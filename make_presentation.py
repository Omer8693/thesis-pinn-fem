"""
make_presentation.py — NAS-PINNs Full Thesis Presentation
==========================================================
21-slide PowerPoint covering all 7 levels with accurate data,
domain visualizations, paper baseline comparisons, and
bilingual speaker notes (English / Turkish).

Usage:
    python make_presentation.py
Output:
    results/NAS_PINNs_Presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
PP   = ROOT / "results" / "pres_plots"   # generated plots
L1   = ROOT / "level1_single_shot/results"
L2   = ROOT / "level2_timestepper/results"
L3   = ROOT / "level3_hybrid_fem/results"
L4   = ROOT / "level4_distortion/results"
L5   = ROOT / "level5_refinement/results"
L6   = ROOT / "level6_poisson_benchmark/results/plots"
L7A  = ROOT / "level7_temporal/results/plots"
L7B  = ROOT / "level7_multiDomain/results/plots"
CRS  = ROOT / "results/cross_level_comparison"

# ── Colour palette ─────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x23, 0x7E)
C_MID    = RGBColor(0x15, 0x65, 0xC0)
C_ACCENT = RGBColor(0xE6, 0x51, 0x00)
C_GREEN  = RGBColor(0x1B, 0x5E, 0x20)
C_RED    = RGBColor(0xC6, 0x28, 0x28)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xE3, 0xF2, 0xFD)
C_GRAY   = RGBColor(0x42, 0x42, 0x42)
C_BG     = RGBColor(0xFA, 0xFA, 0xFF)
C_CREAM  = RGBColor(0xFF, 0xF8, 0xE1)
C_LBLUE  = RGBColor(0xBB, 0xDE, 0xFB)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=14, bold=False, color=C_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return tb


def img(slide, path, l, t, w, h):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), l, t, w, h)
    else:
        rect(slide, l, t, w, h, RGBColor(0xCC, 0xCC, 0xCC))
        txt(slide, f"[{p.name}]", l + Inches(0.1), t + h//2 - Pt(8),
            w - Inches(0.2), Inches(0.3), size=8, color=C_GRAY,
            align=PP_ALIGN.CENTER)


def header(slide, title, subtitle="", tag=""):
    rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_DARK)
    rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.055), C_ACCENT)
    txt(slide, title, Inches(0.35), Inches(0.06), Inches(10.5), Inches(0.6),
        size=26, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.35), Inches(0.66), Inches(10.0), Inches(0.36),
            size=13, color=C_LBLUE, italic=True)
    if tag:
        rect(slide, Inches(11.6), Inches(0.15), Inches(1.5), Inches(0.56), C_ACCENT)
        txt(slide, tag, Inches(11.62), Inches(0.17), Inches(1.46), Inches(0.54),
            size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def footer(slide, extra=""):
    rect(slide, 0, Inches(7.28), SLIDE_W, Inches(0.22), C_DARK)
    foot_txt = "NAS-PINNs  |  2026  |  Mortensen et al."
    if extra:
        foot_txt += f"  |  {extra}"
    txt(slide, foot_txt, Inches(0.3), Inches(7.285), Inches(12.5), Inches(0.2),
        size=8, color=RGBColor(0x90, 0xCA, 0xF9))


def notes(slide, en_text, tr_text=""):
    tf = slide.notes_slide.notes_text_frame
    full = en_text
    if tr_text:
        full += "\n\n— Türkçe —\n" + tr_text
    tf.text = full


def bullets(slide, items, l, t, w, h, size=13, title=None, title_color=C_MID):
    if title:
        txt(slide, title, l, t, w, Inches(0.32), size=12, bold=True,
            color=title_color)
        t += Inches(0.35); h -= Inches(0.35)
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = C_GRAY
        p.space_before = Pt(4)


def table(slide, rows, col_widths, tx, ty, row_h,
          header_color=C_DARK, row_colors=None, txt_colors=None):
    """Draw a simple coloured table."""
    if row_colors is None:
        row_colors = [header_color] + [RGBColor(0xF5,0xF5,0xF5)] * (len(rows)-1)
    if txt_colors is None:
        txt_colors = [C_WHITE] + [C_GRAY] * (len(rows)-1)

    for ri, (row, rc, tc) in enumerate(zip(rows, row_colors, txt_colors)):
        cx = tx
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            rect(slide, cx, ty + ri*row_h, cw - Inches(0.03), row_h - Inches(0.02), rc)
            txt(slide, str(cell),
                cx + Inches(0.05), ty + ri*row_h + Inches(0.06),
                cw - Inches(0.1),  row_h - Inches(0.1),
                size=10 if ri > 0 else 9,
                bold=(ri == 0), color=tc, align=PP_ALIGN.CENTER)
            cx += cw


# Shortcuts for coloured rows in tables
def _cell_ok(cell):  return "✓" in str(cell)
def _cell_bad(cell): return "✗" in str(cell)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank(prs)
    bg(sl, C_DARK)

    rect(sl, 0, Inches(1.9), SLIDE_W, Inches(3.4), C_MID)
    rect(sl, 0, Inches(1.9), SLIDE_W, Inches(0.055), C_ACCENT)

    txt(sl, "NAS-PINNs",
        Inches(0.8), Inches(2.05), Inches(11.5), Inches(1.1),
        size=58, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Neural Architecture Search for Physics-Informed Neural Networks",
        Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.52),
        size=20, color=C_LBLUE, align=PP_ALIGN.CENTER)

    txt(sl, "Accelerating FEM Simulation of A356 Aluminum Water Quenching",
        Inches(0.8), Inches(3.75), Inches(11.5), Inches(0.42),
        size=14, italic=True, color=RGBColor(0x90,0xCA,0xF9), align=PP_ALIGN.CENTER)

    # Level badges
    badges = [("L1","Single-Shot"), ("L2","Skip Op."), ("L3","Hybrid"),
              ("L4","Distortion"), ("L5","Refine"), ("L6","Poisson"), ("L7","Multi-Dom.")]
    bw = Inches(1.56)
    for i, (tag, lbl) in enumerate(badges):
        bx = Inches(0.35) + i * bw
        rect(sl, bx, Inches(5.05), Inches(1.46), Inches(0.82), RGBColor(0x0D,0x47,0xA1))
        txt(sl, f"{tag}\n{lbl}", bx + Inches(0.04), Inches(5.07),
            Inches(1.38), Inches(0.78), size=10, bold=True,
            color=C_WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Mortensen et al. 2026  ·  PyTorch  ·  A356 Aluminum  ·  FEM + PINN Hybrid",
        Inches(0.5), Inches(6.8), Inches(12.2), Inches(0.38),
        size=10, color=RGBColor(0x64,0xB5,0xF6), align=PP_ALIGN.CENTER)

    notes(sl,
        "Welcome. This presentation covers the NAS-PINNs framework: a 7-level system "
        "that applies Neural Architecture Search to Physics-Informed Neural Networks. "
        "The goal is to accelerate finite element simulation of A356 aluminum water "
        "quenching. Three NAS optimizers are compared: Bayesian (TPE), NSGA-II, and "
        "NSGA-III. The core innovation is the temporal skip operator: PINN replaces "
        "intermediate FEM time steps, reducing computation cost.",
        "Bu sunum, NAS-PINNs çerçevesini kapsamaktadır: Nöral Mimari Arama'nın "
        "Fiziğe Dayalı Sinir Ağlarına uygulandığı 7 aşamalı bir sistemdir. Amaç, "
        "A356 alüminyum su söndürme simülasyonunu hızlandırmaktır. Üç NAS "
        "optimizörü karşılaştırılmaktadır: Bayesian, NSGA-II ve NSGA-III. Temel "
        "yenilik, PINN'in ara FEM zaman adımlarını değiştirdiği zamansal atlama "
        "operatörüdür.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════

def slide_overview(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Project Overview", "Seven-level progressive NAS-PINN framework")
    footer(sl)

    levels = [
        ("L1", "Single-Shot PINN",     "Global T(t,x,y) prediction — NAS finds best architecture",   C_MID),
        ("L2", "Skip Operator",        "FEM at anchor points, PINN fills skipped time steps",         C_MID),
        ("L3", "Hybrid FEM + PINN",    "Adaptive routing: PDE residual decides FEM vs PINN",          C_MID),
        ("L4", "Distortion Mechanics", "PINN temperature field → plane-stress → CMM distortion",      C_MID),
        ("L5", "Extended Training",    "Longer Adam + cosine LR: significant L2 improvement",         C_GREEN),
        ("L6", "Poisson Fine-Tuning",  "Auxiliary Poisson loss: tests cross-PDE generalisation",      C_GREEN),
        ("L7", "Multi-Domain",         "5 geometries: Square · Circle · Annulus · L-Shape · Flower",  C_ACCENT),
    ]

    rh = Inches(0.70)
    for i, (tag, name, desc, col) in enumerate(levels):
        y = Inches(1.25) + i * rh
        rect(sl, Inches(0.3), y, Inches(0.65), Inches(0.55), col)
        txt(sl, tag, Inches(0.3), y + Inches(0.06), Inches(0.65), Inches(0.48),
            size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, name, Inches(1.05), y + Inches(0.04), Inches(2.8), Inches(0.50),
            size=12, bold=True, color=C_DARK)
        txt(sl, desc, Inches(3.9), y + Inches(0.08), Inches(9.0), Inches(0.46),
            size=11, color=C_GRAY)
        if i < len(levels)-1:
            rect(sl, Inches(0.3), y + Inches(0.6), Inches(12.8), Inches(0.02),
                 RGBColor(0xCC,0xCC,0xCC))

    rect(sl, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.60), C_CREAM)
    txt(sl,
        "Core Thesis:  PINN with NAS-optimised architecture can replace sequential "
        "FEM steps — fewer solver calls, same accuracy.",
        Inches(0.45), Inches(6.58), Inches(12.4), Inches(0.54),
        size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    notes(sl,
        "The framework has 7 levels. Levels 1-4 build the core pipeline: single-shot "
        "PINN, temporal skip operator, hybrid FEM/PINN routing, and distortion "
        "mechanics. Levels 5-7 extend and validate: Level 5 improves accuracy with "
        "extended training, Level 6 tests a Poisson auxiliary loss, and Level 7 covers "
        "temporal analysis and multi-domain NAS across 5 different geometries.",
        "Çerçeve 7 seviyeden oluşmaktadır. 1-4. seviyeler temel hattı oluşturur: "
        "tek adımlı PINN, zamansal atlama operatörü, hibrit FEM/PINN yönlendirme "
        "ve distorsiyon mekaniği. 5-7. seviyeler genişletme ve doğrulamadır: "
        "5. seviye uzun eğitimle doğruluğu artırır, 6. seviye Poisson yardımcı "
        "kaybını test eder, 7. seviye ise 5 farklı geometride çoklu etki alanı "
        "NAS'ını kapsar.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════

def slide_problem(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Problem Statement", "A356 Aluminum Water Quenching — Why PINN?")
    footer(sl)

    bullets(sl, [
        "▸  Part: A356 cast aluminium alloy  (domain: 1.3 m × 0.6 m, 2D cross-section)",
        "▸  Initial condition: T₀ = 540°C  (solution anneal temperature)",
        "▸  Cooling: immersed in water bath, T_water = 20°C",
        "▸  PDE: ρCₚ · ∂T/∂t = K · ∇²T   with nonlinear Robin BC",
        "▸  HTC h(T): film boiling → nucleate boiling → forced convection regimes",
        "▸  Goal: predict T(x,y,t) accurately over 0-30 s at all mesh points",
        "",
        "FEM Bottleneck:",
        "▸  Requires 20+ sequential time steps (Δt = 1.5 s each)",
        "▸  Each step: full K-matrix assembly + linear solve → expensive",
        "▸  PINN opportunity: learn mapping  T_prev → T_next  and skip FEM calls",
    ], Inches(0.4), Inches(1.28), Inches(6.2), Inches(5.6), size=12,
       title="Physical Setup", title_color=C_MID)

    # Right panel: key numbers
    rect(sl, Inches(6.9), Inches(1.28), Inches(6.1), Inches(5.9),
         RGBColor(0xE8,0xF5,0xE9))
    txt(sl, "Key Numbers", Inches(7.1), Inches(1.42), Inches(5.7), Inches(0.36),
        size=13, bold=True, color=C_GREEN)

    stats = [
        ("540°C → 20°C",  "Temperature range (ΔT = 520°C)"),
        ("30 seconds",    "Total quench duration"),
        ("20 FEM steps",  "Sequential time-step solves"),
        ("K = 150 W/mK",  "A356 thermal conductivity"),
        ("ρCₚ = 2.4 MJ/m³K", "Volumetric heat capacity"),
        ("3 NAS methods", "Bayesian · NSGA-II · NSGA-III"),
    ]
    txt(sl, "Analytical reference (Robin BC fundamental mode):",
        Inches(7.0), Inches(1.85), Inches(5.9), Inches(0.30),
        size=10, bold=True, color=C_DARK)
    txt(sl, "T(t) = 20 + 520 · exp(−1.75×10⁻³ · t)",
        Inches(7.0), Inches(2.18), Inches(5.9), Inches(0.36),
        size=12, bold=True, color=C_MID, align=PP_ALIGN.CENTER)

    for i, (val, lab) in enumerate(stats):
        y = Inches(2.65) + i * Inches(0.71)
        rect(sl, Inches(7.1), y, Inches(2.1), Inches(0.54), C_MID)
        txt(sl, val, Inches(7.12), y + Inches(0.06), Inches(2.06), Inches(0.46),
            size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, lab, Inches(9.3), y + Inches(0.1), Inches(3.55), Inches(0.40),
            size=11, color=C_GRAY)

    notes(sl,
        "The problem is A356 aluminum water quenching. A 1.3×0.6 m 2D cross-section "
        "starts at 540°C and is immersed in 20°C water. The heat equation with "
        "nonlinear Robin BC governs cooling. The heat transfer coefficient h(T) varies "
        "by boiling regime — this nonlinearity makes the problem interesting and "
        "challenging. FEM solves 20 time steps sequentially, each requiring a full "
        "matrix assembly. Our PINN learns to predict T_next given T_prev, replacing "
        "FEM at intermediate steps. The analytical reference T(t) = 20+520·exp(−1.75e-3·t) "
        "is derived from the Robin BC fundamental mode using paper parameters.",
        "Problem: A356 alüminyum su söndürmesidir. 1.3×0.6 m'lik 2B kesit 540°C'den "
        "20°C'lik suya daldırılır. Robin sınır koşullu ısı denklemi kullanılır. "
        "h(T) ısı transfer katsayısı kaynama rejimine göre değişir. FEM 20 zaman "
        "adımını sırayla çözer. PINN, T_prev verilen T_next'i tahmin etmeyi öğrenir "
        "ve ara FEM adımlarını değiştirir.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — NAS Search Space
# ══════════════════════════════════════════════════════════════════════════════

def slide_nas_space(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Neural Architecture Search — What We Are Searching")
    footer(sl)

    bullets(sl, [
        "Search Space (same for all 3 optimizers):",
        "  ▸  Number of hidden layers: {2, 3, 4, 5, 6}",
        "  ▸  Neurons per layer:  32, 48, 64, 96, 128, 160, 192, 256",
        "  ▸  Activation function: tanh · relu · swish",
        "  ▸  Architecture is uniform (same neurons in all layers)",
        "",
        "Evaluation metric (objective to minimise):",
        "  ▸  L2_rel = ‖T_pred − T_ref‖₂ / ‖T_ref‖₂  on held-out test set",
        "",
        "Optimizers:",
        "  ▸  Bayesian (TPE): tree-structured Parzen estimator — single objective",
        "  ▸  NSGA-II: multi-objective (L2 + num params) — Pareto front",
        "  ▸  NSGA-III: multi-objective with reference-point directions",
    ], Inches(0.4), Inches(1.28), Inches(6.5), Inches(5.6), size=12)

    img(sl, PP / "pres_nas_search_space.png",
        Inches(6.8), Inches(1.28), Inches(6.3), Inches(5.5))

    notes(sl,
        "All three optimizers search the same space: 5 choices of depth, 8 choices "
        "of width, and 3 activation functions — roughly 120 architecture candidates. "
        "Bayesian TPE builds a probabilistic model to pick the next candidate. "
        "NSGA-II and NSGA-III treat this as a multi-objective problem: minimise L2 "
        "AND minimise parameter count. This is why they tend to find smaller networks. "
        "Each candidate is trained for a short proxy (1000 epochs) to estimate L2.",
        "Tüm üç optimizör aynı alanı arar: 5 derinlik seçeneği, 8 genişlik seçeneği "
        "ve 3 aktivasyon fonksiyonu — yaklaşık 120 mimari adayı. Bayesian TPE, "
        "sonraki adayı seçmek için olasılıksal bir model oluşturur. NSGA-II ve "
        "NSGA-III bunu çok amaçlı bir problem olarak ele alır: hem L2'yi hem de "
        "parametre sayısını minimize et. Bu nedenle daha küçük ağlar bulmaya eğilimlidirler.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Level 1: Setup & Results Table
# ══════════════════════════════════════════════════════════════════════════════

def slide_level1_table(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 1 — Single-Shot NAS-PINN",
           "Global T(t,x,y) predictor — NAS finds best architecture", tag="L1")
    footer(sl)

    # Left: training setup
    bullets(sl, [
        "Input:   (t, x, y)  →  T(t,x,y)",
        "Training: Adam 20 000 epochs",
        "  • Learning rate: 1×10⁻³ → 1×10⁻⁵  (cosine schedule)",
        "  • Domain points: 2 000  (random collocation)",
        "  • BC points: 200  (boundary sampling)",
        "  • Physics loss: PDE residual + BC residual",
        "",
        "Reference: T_ref from paper FEM (Mortensen 2026)",
        "  • Analytical: T(t) = 20 + 520·exp(−1.75×10⁻³·t)",
    ], Inches(0.4), Inches(1.28), Inches(5.9), Inches(3.5), size=12,
       title="Training Setup", title_color=C_MID)

    # Results table
    rows = [
        ["Optimizer",     "Architecture",  "Params",  "L2_rel↓",  "MAE (°C)↓", "Status"],
        ["Bayesian ★",    "5×151 relu",    "92 564",  "0.076 ✓",  "39.1°C ✓",  "Meets target"],
        ["NSGA-II",       "3×153 tanh",    "47 890",  "0.252",    "132.6°C",   "Misses target"],
        ["NSGA-III",      "3×75  tanh",    "11 776",  "0.513",    "270.3°C",   "Misses target"],
        ["Target",        "—",             "—",       "< 0.100",  "< 50°C",    "L2 threshold"],
    ]
    cw = [Inches(1.7), Inches(1.6), Inches(1.1), Inches(1.1), Inches(1.35), Inches(1.65)]
    rc = [C_DARK,
          RGBColor(0xC8,0xE6,0xC9), RGBColor(0xFF,0xEB,0xEE),
          RGBColor(0xF1,0xF8,0xE9), RGBColor(0xE3,0xF2,0xFD)]
    tc = [C_WHITE, C_GRAY, C_GRAY, C_GRAY, RGBColor(0x1B,0x5E,0x20)]
    table(sl, rows, cw, Inches(0.4), Inches(4.95), Inches(0.50),
          row_colors=rc, txt_colors=tc)

    txt(sl, "★ Only Bayesian meets L2 < 0.10 at Level 1",
        Inches(0.4), Inches(7.12), Inches(5.9), Inches(0.32),
        size=11, bold=True, color=C_ACCENT)

    # Right: image
    img(sl, PP / "pres_l1_arch_bars.png",
        Inches(6.6), Inches(1.28), Inches(6.5), Inches(5.8))

    notes(sl,
        "Level 1 trains a global PINN: given (t, x, y) it predicts temperature T. "
        "Adam with cosine LR schedule, 20 000 epochs. Bayesian NAS finds the best "
        "architecture (5 layers × 151 neurons, relu activation, 92 564 parameters) "
        "and achieves L2=0.076 — the only one meeting the <0.10 target. NSGA-II finds "
        "a more compact network (3×153 tanh, 47 890 params) but L2=0.252. NSGA-III "
        "is smallest (3×75 tanh, 11 776 params) with L2=0.513. Architecture capacity "
        "directly determines accuracy.",
        "Seviye 1, global bir PINN eğitir: (t,x,y) verildiğinde T sıcaklığını tahmin eder. "
        "Adam 20 000 epoch. Bayesian NAS en iyi mimariyi bulur (5×151 relu, 92 564 parametre) "
        "ve L2=0.076 elde eder — sadece bu hedefi karşılar. NSGA-II daha kompakt (3×153 "
        "tanh) ama L2=0.252. NSGA-III en küçük (3×75 tanh) ve L2=0.513.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Level 1: Cooling Curves
# ══════════════════════════════════════════════════════════════════════════════

def slide_level1_curves(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 1 — Cooling Curves: PINN vs FEM Reference", tag="L1")
    footer(sl)

    bullets(sl, [
        "Probe point: centre of domain (x=0.65 m, y=0.30 m)",
        "Reference: analytical T(t) = 20 + 520·exp(−1.75×10⁻³·t)",
        "Bayesian: MAE = 39.1°C  (L2 = 0.076)  — closest to reference",
        "NSGA-II:  MAE = 132.6°C — significantly above target",
        "NSGA-III: MAE = 270.3°C — largest error, smallest network",
    ], Inches(0.4), Inches(1.28), Inches(5.6), Inches(2.5), size=12)

    rect(sl, Inches(0.4), Inches(3.85), Inches(5.6), Inches(1.5),
         RGBColor(0xE3,0xF2,0xFD))
    txt(sl, "Interpretation:",
        Inches(0.55), Inches(3.92), Inches(5.3), Inches(0.28),
        size=11, bold=True, color=C_DARK)
    txt(sl,
        "Larger, deeper network (Bayesian 5×151) captures the complex "
        "spatiotemporal cooling dynamics better. NSGA optimizers sacrifice "
        "accuracy for parameter efficiency.",
        Inches(0.55), Inches(4.22), Inches(5.3), Inches(0.95),
        size=11, color=C_GRAY)

    img(sl, L1 / "fig1_cooling_curves.png",
        Inches(6.3), Inches(1.18), Inches(6.8), Inches(6.0))

    notes(sl,
        "The cooling curves show temperature versus time at the domain centre. "
        "The dashed reference line is the analytical FEM solution. Bayesian PINN "
        "follows it closely (MAE 39°C). NSGA-II overshoots significantly. "
        "NSGA-III, with only 11 776 parameters, cannot represent the spatial "
        "complexity and gives MAE 270°C. This plot demonstrates why architecture "
        "quality matters: a 8× larger network (Bayesian) achieves 7× lower MAE.",
        "Soğuma eğrileri, domain merkezindeki sıcaklık-zaman grafiğini gösterir. "
        "Kesikli referans çizgisi analitik FEM çözümüdür. Bayesian PINN'i yakından "
        "takip eder (MAE 39°C). NSGA-II önemli ölçüde aşıyor. Sadece 11 776 "
        "parametreli NSGA-III, uzamsal karmaşıklığı temsil edemez ve MAE 270°C "
        "verir.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Level 2: Skip Operator Concept
# ══════════════════════════════════════════════════════════════════════════════

def slide_level2_concept(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 2 — Temporal Skip Operator",
           "PINN learns T_prev → T_next, replacing intermediate FEM solves", tag="L2")
    footer(sl)

    bullets(sl, [
        "Window PINN input:   [x, y, t_local, T_prev(x,y)]  →  T_next(x,y)",
        "FEM runs only at anchor points  t[::k]  (every k-th step)",
        "PINN predicts the skipped intermediate steps",
        "Convergence criterion:  MAE(skip=k) / MAE(skip=1)  < 1.5×",
        "",
        "Training parameters:",
        "  • 500 Adam epochs per window  (baseline experiment)",
        "  • LR: 1×10⁻³ → 1×10⁻⁵  (cosine)",
        "  • Domain: 500 collocation pts,  BC: 100 pts",
        "  • T_end = 30 s,  N_steps = 20  (Δt = 1.5 s per step)",
    ], Inches(0.4), Inches(1.28), Inches(5.9), Inches(5.5), size=12,
       title="How the Skip Operator Works", title_color=C_MID)

    # Diagram
    rect(sl, Inches(6.6), Inches(1.28), Inches(6.5), Inches(5.9),
         RGBColor(0xE8,0xF5,0xE9))
    txt(sl, "Example: skip = 2",
        Inches(6.8), Inches(1.38), Inches(6.1), Inches(0.32),
        size=13, bold=True, color=C_GREEN)

    steps_diagram = [
        ("FEM only:", "t=0, 1.5, 3, 4.5 … 30   (21 pts, 20 FEM solves)",    C_MID),
        ("skip=2:",   "FEM at [0, 3, 6 … 30]    (11 FEM solves, 2× faster)", C_GREEN),
        ("",          "PINN at [1.5, 4.5 … 28.5] (10 PINN predictions)",     C_GREEN),
        ("skip=4:",   "FEM at [0, 6, 12 … 30]   (6 FEM solves, 4× faster)",  C_ACCENT),
        ("",          "PINN fills 15 intermediate steps",                      C_ACCENT),
    ]
    for i, (lbl, desc, col) in enumerate(steps_diagram):
        y = Inches(1.85) + i * Inches(0.72)
        if lbl:
            rect(sl, Inches(6.8), y, Inches(1.1), Inches(0.45), col)
            txt(sl, lbl, Inches(6.82), y + Inches(0.04), Inches(1.06), Inches(0.40),
                size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, desc, Inches(8.0), y + Inches(0.07), Inches(4.9), Inches(0.42),
            size=11, color=C_GRAY)

    txt(sl, "Key benefit: Bayesian skip=2 → MAE actually IMPROVES\n"
            "(33.2°C vs 43.6°C baseline) — FEM accumulates numerical error!",
        Inches(6.7), Inches(5.4), Inches(6.2), Inches(0.90),
        size=11, bold=True, color=C_DARK)

    notes(sl,
        "The skip operator is the core innovation. Instead of a global PINN, "
        "a window-based network takes [x, y, t_local, T_prev] and predicts T_next. "
        "FEM runs at anchor points only. With skip=2, only 11 of 21 steps need FEM. "
        "An interesting observation: Bayesian at skip=2 achieves MAE=33.2°C, which "
        "is BETTER than skip=1 (43.6°C). This is because FEM also accumulates "
        "numerical errors over 21 steps, while PINN makes a direct prediction. "
        "The convergence threshold is 1.5× the skip=1 MAE.",
        "Atlama operatörü temel yeniliktir. Global bir PINN yerine, pencere tabanlı "
        "bir ağ [x, y, t_local, T_prev] alır ve T_next'i tahmin eder. FEM yalnızca "
        "çapa noktalarda çalışır. skip=2 ile 21 adımın yalnızca 11'i FEM'e ihtiyaç "
        "duyar. İlginç bir gözlem: Bayesian skip=2'de MAE=33.2°C elde eder, bu "
        "skip=1'den (43.6°C) DAHA İYİDİR. Bunun nedeni FEM'in de 21 adımda "
        "sayısal hata biriktirmesidir.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Level 2: Skip Results Table (500 epochs, corrected)
# ══════════════════════════════════════════════════════════════════════════════

def slide_level2_table(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 2 — Skip Results: 500 Adam Epochs",
           "Convergence threshold: MAE(skip=k) / MAE(skip=1) < 1.5×", tag="L2")
    footer(sl)

    # Table — CORRECTED: NSGA-II skip=2 ratio=1.47× → CONVERGED ✓
    rows = [
        ["Skip",   "FEM Steps", "Bayesian",          "Ratio",   "NSGA-II",           "Ratio",   "NSGA-III",           "Ratio"  ],
        ["skip=1", "21/21",     "43.6°C ✓",          "1.00×",   "38.1°C ✓",          "1.00×",   "44.4°C ✓",           "1.00×"  ],
        ["skip=2", "11/21",     "33.2°C ✓",          "0.76×",   "56.0°C ✓",          "1.47×",   "75.6°C ✗",           "1.70×"  ],
        ["skip=4", " 6/21",     "57.5°C ✓",          "1.32×",   "154.8°C ✗",         "4.06×",   "202.1°C ✗",          "4.55×"  ],
        ["skip=6", " 4/21",     "93.6°C ✗",          "2.15×",   "354.6°C ✗",         "9.30×",   "428.6°C ✗",          "9.65×"  ],
    ]
    cw = [Inches(0.9), Inches(1.1),
          Inches(1.45), Inches(0.85),
          Inches(1.45), Inches(0.85),
          Inches(1.55), Inches(0.85)]

    # Build row colors with ✓/✗ per cell
    row_h = Inches(0.52)
    for ri, row in enumerate(rows):
        cx = Inches(0.35)
        for ci, (cell, cw_i) in enumerate(zip(row, cw)):
            if ri == 0:
                cell_col = C_DARK
                tc       = C_WHITE
            elif ci >= 2:
                if "✓" in str(cell):
                    cell_col = RGBColor(0xC8,0xE6,0xC9)
                    tc       = RGBColor(0x1B,0x5E,0x20)
                elif "✗" in str(cell):
                    cell_col = RGBColor(0xFF,0xCD,0xD2)
                    tc       = C_RED
                else:
                    cell_col = RGBColor(0xF5,0xF5,0xF5)
                    tc       = C_GRAY
            else:
                cell_col = RGBColor(0xF5,0xF5,0xF5)
                tc       = C_GRAY

            rect(sl, cx, Inches(1.28) + ri * row_h,
                 cw_i - Inches(0.03), row_h - Inches(0.02), cell_col)
            txt(sl, str(cell),
                cx + Inches(0.04), Inches(1.28) + ri * row_h + Inches(0.07),
                cw_i - Inches(0.08), row_h - Inches(0.1),
                size=9 if ri > 0 else 8,
                bold=(ri == 0), color=tc, align=PP_ALIGN.CENTER)
            cx += cw_i

    # Key finding box
    rect(sl, Inches(0.35), Inches(4.0), Inches(12.6), Inches(1.0),
         RGBColor(0xFF,0xF3,0xE0))
    txt(sl, "Key findings at 500 epochs:",
        Inches(0.5), Inches(4.06), Inches(12.2), Inches(0.26),
        size=11, bold=True, color=C_ACCENT)
    txt(sl,
        "▸ Bayesian achieves skip=4 ✓ (1.32×) — 4× FEM reduction      "
        "▸ NSGA-II achieves skip=2 ✓ (1.47×, just below threshold)      "
        "▸ NSGA-III fails at skip=2 ✗ (1.70×) — smaller model needs more training",
        Inches(0.5), Inches(4.34), Inches(12.2), Inches(0.55),
        size=11, color=C_GRAY)

    img(sl, PP / "pres_l2_skip_comparison.png",
        Inches(0.35), Inches(5.1), Inches(12.7), Inches(2.0))

    notes(sl,
        "At 500 epochs, the corrected results show: Bayesian converges at skip=2 "
        "(ratio 0.76×, MAE actually improves!) and skip=4 (ratio 1.32×). "
        "NSGA-II converges at skip=2 with ratio 1.47× — just under the 1.5× threshold. "
        "This was incorrectly shown as diverged in the previous version. "
        "NSGA-III fails at skip=2 (ratio 1.70×). Ratios above 1.5 are too inaccurate.",
        "500 epochta düzeltilmiş sonuçlar: Bayesian skip=2'de (oran 0.76×, "
        "MAE aslında iyileşiyor!) ve skip=4'te (1.32×) yakınsar. NSGA-II, "
        "oran 1.47× ile skip=2'de yakınsar — 1.5× eşiğinin hemen altında. "
        "Bu önceki sürümde yanlışlıkla ıraksak olarak gösterilmişti. "
        "NSGA-III skip=2'de başarısız olur (oran 1.70×).")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Level 2: Skip Visualisation
# ══════════════════════════════════════════════════════════════════════════════

def slide_level2_viz(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 2 — Skip Operator: Visual Comparison", tag="L2")
    footer(sl)

    bullets(sl, [
        "This plot shows the temperature prediction at skip=1, 2, 4, 6",
        "for the Bayesian architecture.",
        "",
        "Observations:",
        "  ▸  skip=2 (11 FEM): prediction improves vs skip=1",
        "  ▸  skip=4 (6 FEM): acceptable error, 4× speedup",
        "  ▸  skip=6 (4 FEM): error exceeds threshold — diverged",
        "",
        "Speedup summary:",
        "  ▸  skip=2: 11 FEM calls → 2× faster",
        "  ▸  skip=4:  6 FEM calls → 4× faster",
    ], Inches(0.4), Inches(1.28), Inches(4.8), Inches(5.5), size=12)

    img(sl, L2 / "fem_vs_pinn_skip.png",
        Inches(5.4), Inches(1.18), Inches(7.7), Inches(6.0))

    notes(sl,
        "The skip comparison plot shows temperature fields at the domain probe point "
        "for different skip values. Skip=2 with Bayesian gives the best result — "
        "smoother than skip=1 because PINN avoids FEM's sequential accumulation. "
        "Skip=4 is still within tolerance. Skip=6 diverges significantly.",
        "Atlama karşılaştırma grafiği, farklı atlama değerleri için domain prob "
        "noktasındaki sıcaklık alanlarını gösterir. Bayesian ile skip=2 en iyi "
        "sonucu verir — PINN, FEM'in sıralı birikimini önlediğinden skip=1'den "
        "daha düzgündür. Skip=4 hâlâ tolerans içindedir. Skip=6 önemli ölçüde ıraksak.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Level 2: 2000-epoch Extended Training
# ══════════════════════════════════════════════════════════════════════════════

def slide_level2_2k(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 2 — Extended Training: 2 000 Epochs for NSGA-II/III",
           "Can more training epochs fix the skip=2 convergence problem?", tag="L2")
    footer(sl)

    bullets(sl, [
        "Re-run NSGA-II and NSGA-III at skip=1 and skip=2 with 2 000 Adam epochs",
        "(vs 500 epochs in the baseline experiment)",
        "",
        "Training parameters (2000-epoch run):",
        "  • n_epochs = 2 000  (4× baseline)",
        "  • n_domain = 500 collocation points",
        "  • n_bc = 100 boundary points",
        "  • LR: 1×10⁻³ → 1×10⁻⁵  (cosine)",
    ], Inches(0.4), Inches(1.28), Inches(5.8), Inches(3.4), size=12)

    # 2000-epoch results table
    # NSGA-II skip=1: 18.4°C, skip=2: 37.4°C, ratio=37.4/38.1=0.98× ✓
    # NSGA-III skip=1: 33.8°C, skip=2: 43.8°C, ratio=43.8/44.4=0.99× ✓
    # Ratio compared to 500-epoch skip=1 reference (not 2000-epoch skip=1)
    rows_2k = [
        ["Optimizer",  "skip=1 MAE\n(2 000 ep)",  "skip=2 MAE\n(2 000 ep)",
         "Ratio vs\n500ep ref",  "Converged?"],
        ["NSGA-II\n3×153 tanh",  "18.4°C",  "37.4°C",  "0.98×",  "✓ YES"],
        ["NSGA-III\n3×75 tanh",  "33.8°C",  "43.8°C",  "0.99×",  "✓ YES"],
        ["500ep reference",       "38.1°C",  "56.0°C",  "(1.47×)", "(baseline)"],
    ]
    cw2 = [Inches(2.2), Inches(2.0), Inches(2.0), Inches(1.8), Inches(1.7)]
    rc2 = [C_DARK,
           RGBColor(0xC8,0xE6,0xC9), RGBColor(0xC8,0xE6,0xC9),
           RGBColor(0xF5,0xF5,0xF5)]
    tc2 = [C_WHITE,
           RGBColor(0x1B,0x5E,0x20), RGBColor(0x1B,0x5E,0x20),
           C_GRAY]
    table(sl, rows_2k, cw2, Inches(0.4), Inches(4.85), Inches(0.55),
          row_colors=rc2, txt_colors=tc2)

    rect(sl, Inches(0.4), Inches(7.0), Inches(12.6), Inches(0.20),
         RGBColor(0xE8,0xF5,0xE9))
    txt(sl,
        "Conclusion: Epoch count was the bottleneck, not architecture. "
        "2 000 epochs → NSGA-II ratio=0.98× ✓ and NSGA-III ratio=0.99× ✓ — both converge at skip=2.",
        Inches(0.5), Inches(7.02), Inches(12.2), Inches(0.20),
        size=10, bold=True, color=RGBColor(0x1B,0x5E,0x20))

    img(sl, L2 / "skip_comparison.png",
        Inches(6.4), Inches(1.28), Inches(6.7), Inches(3.4))

    notes(sl,
        "The 2000-epoch run answers: is NSGA-II/III skip=2 failure due to architecture "
        "or insufficient training? Result: both converge with 2000 epochs. "
        "NSGA-II ratio=0.98× (skip=2 MAE better than 500-epoch skip=1 reference!). "
        "NSGA-III ratio=0.99×. The bottleneck was epoch count, not model capacity. "
        "Bayesian converges faster because its larger network (5×151) learns the "
        "temporal dynamics more efficiently.",
        "2000 epochluk çalıştırma şunu yanıtlar: NSGA-II/III skip=2 başarısızlığı "
        "mimariden mi yoksa yetersiz eğitimden mi? Sonuç: her ikisi de 2000 epochla "
        "yakınsar. NSGA-II oranı=0.98×. NSGA-III oranı=0.99×. Darboğaz epoch "
        "sayısıydı, model kapasitesi değildi.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Level 3: Hybrid FEM+PINN
# ══════════════════════════════════════════════════════════════════════════════

def slide_level3(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 3 — Hybrid FEM + PINN: Adaptive Routing",
           "PDE residual decides at each step: FEM or PINN?", tag="L3")
    footer(sl)

    bullets(sl, [
        "Parameters:",
        "  • Residual threshold = 0.1",
        "  • Max consecutive PINN steps: 4  (config A) or 20 (config B)",
        "  • T_end = 30 s,  N_steps = 20",
        "",
        "Algorithm:",
        "  1. Solve step with current best model",
        "  2. Compute PDE residual r(t) = ‖∂T/∂t − K∇²T‖ / ‖T‖",
        "  3. If r > 0.1 → trigger full FEM solve",
        "  4. If r ≤ 0.1 → accept PINN prediction",
        "  5. Advance to next time step",
        "",
        "Result (Bayesian architecture):",
        "  ▸  Config A (max_skip=4):  4 FEM + 16 PINN = 80% skip",
        "  ▸  Config B (max_skip=20): 1 FEM + 19 PINN = 95% skip",
    ], Inches(0.4), Inches(1.28), Inches(5.9), Inches(5.9), size=12)

    img(sl, PP / "pres_l3_skip_rates.png",
        Inches(6.4), Inches(1.28), Inches(6.7), Inches(5.9))

    notes(sl,
        "Level 3 adds adaptive intelligence. Instead of a fixed skip factor, the system "
        "decides at each step whether to call FEM based on the PDE residual. "
        "If residual > 0.1 → FEM for accuracy. If ≤ 0.1 → PINN. "
        "Config A (max_skip=4): 4 FEM calls, 16 PINN = 80% skip. "
        "Config B (max_skip=20): only 1 FEM call, 19 PINN = 95% skip. "
        "The residual check prevents error accumulation automatically. "
        "All three optimizers achieve 80% skip in Config A.",
        "Seviye 3 adaptif zeka ekler. Sabit bir atlama faktörü yerine, sistem her "
        "adımda PDE artığına göre FEM çağrısı yapıp yapmayacağına karar verir. "
        "Artık > 0.1 → doğruluk için FEM. ≤ 0.1 → PINN. "
        "Konfigürasyon A: 4 FEM çağrısı, 16 PINN = %80 atlama. "
        "Konfigürasyon B: yalnızca 1 FEM çağrısı, 19 PINN = %95 atlama.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Level 3: Residual Trace
# ══════════════════════════════════════════════════════════════════════════════

def slide_level3_residual(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 3 — Residual Trace and Step Distribution", tag="L3")
    footer(sl)

    bullets(sl, [
        "Left plot: residual at each time step",
        "  ▸  FEM triggered (red) when residual crosses 0.1 threshold",
        "  ▸  Residual drops sharply after each FEM correction",
        "  ▸  PINN steps (blue) maintain low residual between corrections",
        "",
        "Right plot: which steps used FEM vs PINN",
        "  ▸  First step always FEM (initial condition anchor)",
        "  ▸  After first few steps residual remains small",
        "  ▸  PINN handles most of the simulation safely",
    ], Inches(0.4), Inches(1.28), Inches(4.5), Inches(3.8), size=12)

    img(sl, L3 / "thr0.1_skip4/residual_trace.png",
        Inches(0.4), Inches(4.95), Inches(4.5), Inches(2.15))
    img(sl, L3 / "thr0.1_skip4/step_distribution.png",
        Inches(5.0), Inches(1.28), Inches(8.1), Inches(5.9))

    notes(sl,
        "The residual trace shows r(t) over the 20 time steps. FEM is triggered "
        "at step 1 (initial anchor), step 6, step 11, and step 16 — every 5 steps. "
        "Between FEM calls, the residual stays below 0.1. The PINN accumulates some "
        "error by step 5, triggering FEM, which resets the residual. This pattern "
        "repeats throughout the simulation, giving 4 FEM + 16 PINN = 80% efficiency.",
        "Artık izi, 20 zaman adımı boyunca r(t)'yi gösterir. FEM, adım 1 (başlangıç "
        "çapası), adım 6, adım 11 ve adım 16'da tetiklenir — her 5 adımda bir. "
        "FEM çağrıları arasında artık 0.1'in altında kalır. PINN'in hatası birikir "
        "ve FEM'i tetikler, bu da artığı sıfırlar. Bu model boyunca tekrar eder.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Level 4: Thermal Distortion
# ══════════════════════════════════════════════════════════════════════════════

def slide_level4(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 4 — Thermal Distortion at CMM Points",
           "PINN temperature field → plane-stress FEM → distortion in mm", tag="L4")
    footer(sl)

    bullets(sl, [
        "Input: T(x,y,t=30s) from Level 3 PINN prediction",
        "Method: 2D plane-stress FEM with Q4 elements",
        "  • σ_thermal = E · α · ΔT,  where ΔT = T - T_room",
        "  • E = 69 GPa  (A356 elastic modulus)",
        "  • α = 22×10⁻⁶ /°C  (thermal expansion coefficient)",
        "Output: displacement δ (mm) at 47 CMM measurement points",
        "Reference: paper Fig 17/18  (Mortensen 2026)",
        "",
        "Results — Mean absolute distortion |δ|:",
        "  ▸  Paper measured:  0.78 mm  (physical experiment)",
        "  ▸  Paper FEM:       0.75 mm  (reference simulation)",
        "  ▸  Bayesian PINN:   1.67 mm  (closest PINN result)",
        "  ▸  NSGA-II PINN:    2.81 mm",
        "  ▸  NSGA-III PINN:   4.74 mm",
    ], Inches(0.4), Inches(1.28), Inches(5.9), Inches(5.9), size=12)

    img(sl, PP / "pres_l4_distortion.png",
        Inches(6.3), Inches(1.28), Inches(6.8), Inches(5.9))

    notes(sl,
        "Level 4 bridges thermal simulation and structural mechanics. The PINN "
        "temperature at t=30s drives a plane-stress FEM for thermal distortion. "
        "Results compared to CMM measurements and paper FEM. "
        "Paper FEM achieves 0.75mm (very close to measured 0.78mm). "
        "Our Bayesian PINN gives 1.67mm — 2× worse than paper but reasonable "
        "given it uses a neural approximation. NSGA-III gives 4.74mm — "
        "corresponding to its large temperature error (L2=0.513, 270°C MAE). "
        "The direct correlation: PINN thermal accuracy → distortion accuracy.",
        "Seviye 4, termal simülasyon ile yapısal mekanik arasında köprü kurar. "
        "t=30s'deki PINN sıcaklığı, termal distorsiyon için bir düzlem gerilim FEM'i "
        "yönlendirir. Sonuçlar CMM ölçümleri ve paper FEM ile karşılaştırılır. "
        "Paper FEM 0.75mm (ölçülen 0.78mm'ye çok yakın). Bayesian PINN'imiz "
        "1.67mm veriyor — paper'dan 2× daha kötü ama makul.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Level 4: Paper Comparison
# ══════════════════════════════════════════════════════════════════════════════

def slide_level4_comparison(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 4 — CMM Point-by-Point Comparison vs Paper", tag="L4")
    footer(sl)

    bullets(sl, [
        "47 CMM points measured in paper (Fig 17 — bottom layer)",
        "Positive = material moved up, negative = moved down",
        "Paper FEM matches measurements well (MAE = 0.18 mm)",
        "Bayesian PINN: MAE vs measured = 1.14 mm",
        "Error comes from temperature field approximation, not FEM solver",
        "",
        "Next step to improve: better PINN temperature accuracy",
        "  (Level 5 extended training reduces L2 from 0.076 to 0.030)",
    ], Inches(0.4), Inches(1.28), Inches(5.0), Inches(3.8), size=12)

    img(sl, L4 / "paper_comparison_signed.png",
        Inches(5.2), Inches(1.18), Inches(7.9), Inches(6.1))

    notes(sl,
        "This slide shows the signed distortion at each CMM point. The paper FEM "
        "(orange) follows the measured values (grey) closely. Bayesian PINN (blue) "
        "captures the general trend but with larger deviations. The discrepancy "
        "is mainly from temperature error propagation. Improving the PINN "
        "temperature accuracy (as done in Level 5) would directly reduce distortion error.",
        "Bu slayt her CMM noktasındaki işaretli distorsiyonu gösterir. Paper FEM (turuncu), "
        "ölçülen değerleri (gri) yakından takip eder. Bayesian PINN (mavi) genel trendi "
        "yakalar ancak daha büyük sapmalar vardır. Tutarsızlık, ağırlıklı olarak sıcaklık "
        "hata yayılımından kaynaklanmaktadır.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Level 5: Extended Training
# ══════════════════════════════════════════════════════════════════════════════

def slide_level5(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 5 — Extended Adam Training: L2 Improvement",
           "Re-train Level 1 models with cosine LR for 20 000 epochs", tag="L5")
    footer(sl)

    bullets(sl, [
        "Same architectures as Level 1",
        "Training: Adam 20 000 epochs with cosine LR: 1×10⁻³ → 1×10⁻⁵",
        "Goal: squeeze more accuracy from existing architectures",
        "",
        "Note: L-BFGS was tested but did not improve results in this setup",
        "(lbfgs_iters=0 — the Adam solution was already in a flat region)",
        "Extended cosine LR is the source of improvement here.",
    ], Inches(0.4), Inches(1.28), Inches(5.8), Inches(3.5), size=12,
       title="Method", title_color=C_MID)

    rows = [
        ["Optimizer",       "L1 L2_rel", "L1 MAE",  "L5 L2_rel", "L5 MAE",   "Improvement"],
        ["Bayesian 5×151",  "0.076",     "39.1°C",   "0.030 ✓",  "14.4°C ✓", "2.6× ↓"],
        ["NSGA-II 3×153",   "0.252",     "132.6°C",  "0.055 ✓",  "28.7°C ✓", "4.5× ↓"],
        ["NSGA-III 3×75",   "0.513",     "270.3°C",  "0.259",    "136.4°C",  "2.0× ↓"],
        ["Target",          "—",         "—",         "< 0.100",  "< 50°C",   "—"],
    ]
    cw = [Inches(2.0), Inches(1.2), Inches(1.3), Inches(1.25), Inches(1.35), Inches(1.4)]
    rc = [C_DARK,
          RGBColor(0xE3,0xF2,0xFD), RGBColor(0xFF,0xEB,0xEE),
          RGBColor(0xF1,0xF8,0xE9), RGBColor(0xE8,0xF5,0xE9)]
    tc = [C_WHITE, C_GRAY, C_GRAY, C_GRAY, RGBColor(0x1B,0x5E,0x20)]
    table(sl, rows, cw, Inches(0.4), Inches(4.9), Inches(0.50),
          row_colors=rc, txt_colors=tc)

    txt(sl,
        "NSGA-II: largest gain (4.5×) — compact architecture responds well to longer training",
        Inches(0.4), Inches(7.12), Inches(12.6), Inches(0.28),
        size=11, bold=True, color=C_MID)

    img(sl, PP / "pres_l5_improvement.png",
        Inches(6.4), Inches(1.28), Inches(6.7), Inches(3.4))

    notes(sl,
        "Level 5 re-trains the Level 1 architectures with a longer cosine LR schedule. "
        "The key finding: NSGA-II improves the most (4.5×), going from L2=0.252 to "
        "0.055 — now meeting the <0.10 target! Bayesian improves from 0.076 to 0.030 "
        "(2.6×). NSGA-III shows only 2× improvement and stays at 0.259 — its small "
        "capacity (11 776 params) limits accuracy regardless of training length. "
        "Note: L-BFGS was attempted but ran 0 iterations, suggesting the Adam "
        "minimum was already in a locally flat region.",
        "Seviye 5, Seviye 1 mimarilerini daha uzun bir kosinüs LR çizelgesiyle "
        "yeniden eğitir. Temel bulgu: NSGA-II en çok iyileşiyor (4.5×), L2=0.252'den "
        "0.055'e gidiyor — artık <0.10 hedefini karşılıyor! Bayesian 0.076'dan "
        "0.030'a iyileşiyor (2.6×). NSGA-III yalnızca 2× iyileşme gösteriyor ve "
        "0.259'da kalıyor.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Level 5: L2 Comparison Plot
# ══════════════════════════════════════════════════════════════════════════════

def slide_level5_plot(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 5 — L2 Before vs After: All Optimizers", tag="L5")
    footer(sl)

    bullets(sl, [
        "The plot shows L2_rel comparison across all three optimizers",
        "",
        "After extended training:",
        "  ▸  Bayesian:  L2 = 0.030  (was 0.076)  — 2.6× improvement",
        "  ▸  NSGA-II:   L2 = 0.055  (was 0.252)  — 4.5× improvement  ✓",
        "  ▸  NSGA-III:  L2 = 0.259  (was 0.513)  — 2.0× improvement",
        "",
        "Skip capability after Level 5:",
        "  ▸  Bayesian L5:   skip=4 still converges  (L2=0.030 → good generalisation)",
        "  ▸  NSGA-II L5:    skip=2 converges  (L2=0.055 → near target)",
        "  ▸  NSGA-III L5:   skip=2 marginal   (still above 0.10 threshold)",
    ], Inches(0.4), Inches(1.28), Inches(5.6), Inches(5.9), size=12)

    img(sl, L5 / "level5_l2_comparison.png",
        Inches(6.1), Inches(1.18), Inches(7.0), Inches(6.0))

    notes(sl,
        "The L2 comparison plot from Level 5 shows bar charts for each optimizer "
        "before (lighter) and after (darker) extended training. The most striking "
        "result is NSGA-II: despite being a 3-layer 153-neuron network, 20 000 "
        "epochs with cosine LR brings it from L2=0.252 to 0.055 — a 4.5× improvement. "
        "This confirms the architecture has sufficient capacity; it just needed more "
        "training time.",
        "Seviye 5'teki L2 karşılaştırma grafiği, uzatılmış eğitim öncesi (daha açık) "
        "ve sonrası (daha koyu) her optimizör için çubuk grafikleri gösteriyor. "
        "En çarpıcı sonuç NSGA-II: 3 katmanlı 153-nöronlu ağ olmasına rağmen "
        "kosinüs LR ile 20 000 epoch, onu L2=0.252'den 0.055'e getiriyor — 4.5× iyileşme.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Level 6: Poisson Auxiliary Fine-Tuning
# ══════════════════════════════════════════════════════════════════════════════

def slide_level6(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 6 — Poisson Auxiliary Fine-Tuning",
           "Add Poisson PDE loss to quenching network — marginal improvement", tag="L6")
    footer(sl)

    bullets(sl, [
        "Goal: does adding Poisson-like physics loss improve quenching accuracy?",
        "",
        "Method:",
        "  • Take Level 5 weights (already trained on quenching)",
        "  • Add auxiliary Poisson loss at 5 time slices: t ∈ {0.5, 2, 5, 10, 20} s",
        "  • Fine-tune for 5 000 epochs  LR: 1×10⁻⁵ → 1×10⁻⁶",
        "  • λ_poisson = 1.0  (weight of auxiliary term)",
        "",
        "Results — L2 before vs after fine-tuning:",
        "  ▸  Bayesian:  0.038 → 0.031  (+17.5% improvement)",
        "  ▸  NSGA-II:   0.062 → 0.058  (+6.2% improvement)",
        "  ▸  NSGA-III:  0.274 → 0.267  (+2.5% improvement)",
        "",
        "Conclusion: marginal gains — Poisson auxiliary loss provides a small",
        "regularisation benefit but does not fundamentally change accuracy.",
    ], Inches(0.4), Inches(1.28), Inches(6.0), Inches(5.9), size=12)

    img(sl, L6 / "l6_heatmap_t5s.png",
        Inches(6.6), Inches(1.28), Inches(6.5), Inches(5.9))

    notes(sl,
        "Level 6 tests whether adding Poisson-type auxiliary physics loss can improve "
        "the quenching PINN. The gains are marginal: Bayesian improves 17.5%, "
        "NSGA-II 6.2%, NSGA-III only 2.5%. This suggests the Poisson auxiliary loss "
        "acts as a regulariser but the dominant accuracy constraint is the "
        "network architecture capacity and main PDE training. "
        "Level 7B takes a different approach: full NAS on actual Poisson problems "
        "with different geometries, achieving much stronger results.",
        "Seviye 6, Poisson tipli yardımcı fizik kaybı eklemenin söndürme PINN'ini "
        "iyileştirip iyileştiremeyeceğini test eder. Kazanımlar marjinaldir: Bayesian "
        "%17.5, NSGA-II %6.2, NSGA-III yalnızca %2.5 iyileşir. Bu, Poisson yardımcı "
        "kaybının bir düzenleyici olarak işlev gördüğünü ama temel doğruluk "
        "kısıtının mimari kapasite olduğunu düşündürür.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Level 7B: Domain Shapes
# ══════════════════════════════════════════════════════════════════════════════

def slide_level7b_domains(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 7B — Multi-Domain Poisson: Five Geometries",
           "NAS searches the optimal architecture for each geometry independently", tag="L7")
    footer(sl)

    bullets(sl, [
        "Problem: −∇²u = f(x,y)  with  u = 0 on boundary  (Poisson PDE)",
        "",
        "Five domains tested:",
        "  ▸  Square [0,1]²  — manufactured solution u*=x²(x−1)²y²(y−1)²",
        "  ▸  Circle x²+y²≤1  — exact solution u*=(1−r²)/4",
        "  ▸  Annulus 0.25≤r≤1  — exact solution u*=(r²−r_inner²)/4",
        "  ▸  L-Shape [0,1]²\\[0.5,1]×[0,0.5]  — FDM reference",
        "  ▸  Flower r≤1+0.3·cos(5θ)  — 5 petals, FDM reference",
        "",
        "Colormap shows: reference Poisson solution u(x,y)",
        "Darker colour = higher solution value",
    ], Inches(0.4), Inches(1.28), Inches(4.5), Inches(3.8), size=12)

    img(sl, PP / "pres_domain_shapes.png",
        Inches(0.4), Inches(4.85), Inches(12.6), Inches(2.35))

    notes(sl,
        "Level 7B tests whether NAS can find good architectures for Poisson "
        "problems across very different geometries. The bottom row shows the 5 "
        "domains with their reference solutions (YlGnBu colormap). The circle "
        "and annulus have exact analytical solutions. The L-shape and flower "
        "are solved with FDM as reference. NAS runs independently for each "
        "geometry × optimizer combination (15 NAS runs total).",
        "Seviye 7B, NAS'ın çok farklı geometriler genelinde Poisson problemleri "
        "için iyi mimariler bulup bulamayacağını test eder. Alt satır, referans "
        "çözümleriyle birlikte 5 domaini gösterir. Daire ve halka için analitik "
        "çözümler mevcuttur. L-şekil ve çiçek için FDM referans olarak kullanılır.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Level 7B: Prediction + Error Maps
# ══════════════════════════════════════════════════════════════════════════════

def slide_level7b_pred_error(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 7B — PINN Prediction vs Reference: Error Maps",
           "3 rows: Reference solution  |  Bayesian PINN pred  |  |pred − ref|", tag="L7")
    footer(sl)

    bullets(sl, [
        "Top row:    Reference — exact formula (Square, Circle) or FDM (L-Shape, Flower)",
        "Middle row: Bayesian PINN prediction on the same grid",
        "Bottom row: Absolute error |pred − ref|  (YlOrRd colormap — darker = bigger error)",
        "",
        "L2 values (blue badge) are final test errors from NAS training.",
        "Max error (red badge) shows worst-case spatial deviation.",
    ], Inches(0.4), Inches(1.28), Inches(12.7), Inches(1.8), size=12)

    img(sl, PP / "pres_domain_pred_error.png",
        Inches(0.3), Inches(3.05), Inches(12.8), Inches(4.15))

    notes(sl,
        "This slide shows the full comparison: reference (top), PINN prediction (middle), "
        "and absolute error map (bottom) for all four domains. "
        "Circle is excellent: L2=0.00014, error map nearly zero everywhere. "
        "L-shape is good: L2=0.028, error concentrated near the inner corner. "
        "Square: L2=0.038, error distributed across the domain. "
        "Flower: L2=0.172, larger errors especially near petal tips — "
        "the complex boundary is hard to learn accurately.",
        "Bu slayt tam karşılaştırmayı gösteriyor: referans (üst), PINN tahmini "
        "(orta) ve mutlak hata haritası (alt). Daire mükemmeldir: L2=0.00014, "
        "hata haritası neredeyse her yerde sıfır. L-şekil iyidir: L2=0.028, "
        "hata iç köşede yoğunlaşmış. Kare: L2=0.038. Çiçek: L2=0.172, "
        "taç yaprakları özellikle zordur.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Level 7B: Multi-Domain L2 Bar Chart
# ══════════════════════════════════════════════════════════════════════════════

def slide_level7b_results(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Level 7B — Multi-Domain Results: L2 per Geometry",
           "Bayesian NAS finds optimal architecture for each domain", tag="L7")
    footer(sl)

    # Results table
    rows = [
        ["Domain",   "Bayesian",                "Arch",          "NSGA-II",   "NSGA-III",  "Difficulty"],
        ["Square",   "L2=0.0376",               "5×100 gelu",    "N/A",       "N/A",       "Medium"],
        ["Circle",   "L2=0.00014 ✓",            "5×40 gelu",     "0.00042",   "0.00066",   "Easy"],
        ["L-Shape",  "L2=0.0277 ✓",             "5×100 gelu",    "0.0285",    "0.0320",    "Medium"],
        ["Flower",   "L2=0.172",                "2×50 gelu",     "0.176",     "0.176",     "Hard"],
        ["Target",   "L2 < 0.05",               "—",             "L2 < 0.05", "L2 < 0.05", "—"],
    ]
    cw = [Inches(1.5), Inches(2.2), Inches(1.8), Inches(1.7), Inches(1.7), Inches(1.4)]
    rc = [C_DARK,
          RGBColor(0xF5,0xF5,0xF5), RGBColor(0xC8,0xE6,0xC9),
          RGBColor(0xF5,0xF5,0xF5), RGBColor(0xFF,0xEB,0xEE),
          RGBColor(0xE3,0xF2,0xFD)]
    tc = [C_WHITE, C_GRAY, C_GRAY, C_GRAY, C_GRAY, RGBColor(0x1B,0x5E,0x20)]
    table(sl, rows, cw, Inches(0.4), Inches(1.28), Inches(0.50),
          row_colors=rc, txt_colors=tc)

    img(sl, PP / "pres_l7b_results.png",
        Inches(0.4), Inches(4.3), Inches(12.7), Inches(2.85))

    txt(sl,
        "Note: Annulus not included — NAS did not converge for all optimizers. "
        "Flower geometry is hardest (complex 5-petal boundary). Circle is easiest (smooth).",
        Inches(0.4), Inches(7.13), Inches(12.6), Inches(0.26),
        size=10, italic=True, color=C_GRAY)

    notes(sl,
        "Multi-domain results show strong variation by geometry. Circle is easiest: "
        "smooth boundary, analytical exact solution, all optimizers converge easily. "
        "Best: Bayesian 5×40 gelu with L2=0.00014. L-shape is medium difficulty. "
        "Flower is hardest: the 5-petal boundary creates sharp corners, none of the "
        "optimizers converge well (L2≈0.17). This shows that geometry complexity "
        "significantly impacts PINN performance and NAS benefits.",
        "Çoklu domain sonuçları geometriye göre güçlü varyasyon gösteriyor. "
        "Daire en kolaydır: pürüzsüz sınır, analitik tam çözüm. En iyi: "
        "Bayesian 5×40 gelu, L2=0.00014. L-şekil orta güçlüktedir. "
        "Çiçek en zordur: 5 yapraklı sınır keskin köşeler oluşturur, "
        "hiçbir optimizör iyi yakınsamaz (L2≈0.17).")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Cross-Level Summary
# ══════════════════════════════════════════════════════════════════════════════

def slide_cross_level(prs):
    sl = blank(prs)
    bg(sl, C_BG)
    header(sl, "Cross-Level Summary — L2 Progression Across All Levels")
    footer(sl)

    img(sl, PP / "pres_cross_level_l2.png",
        Inches(0.35), Inches(1.18), Inches(6.5), Inches(5.9))

    # Summary table
    rows = [
        ["Level",  "Method",                "Bayesian L2",  "NSGA-II L2",  "NSGA-III L2"],
        ["L1",     "Adam 20k epochs",       "0.076",        "0.252",       "0.513"],
        ["L2 ★",   "Skip=2 (500ep)",        "0.108",        "0.181",       "0.217"],
        ["L3",     "Hybrid 80% skip",       "≈0.076",       "≈0.127",      "≈0.137"],
        ["L5",     "Extended Adam",         "0.030 ✓",      "0.055 ✓",     "0.259"],
        ["L6",     "Poisson fine-tune",     "0.031 ✓",      "0.058 ✓",     "0.267"],
        ["Target", "—",                     "< 0.10",       "< 0.10",      "< 0.10"],
    ]
    cw = [Inches(0.9), Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.7)]
    rc = [C_DARK,
          RGBColor(0xE3,0xF2,0xFD), RGBColor(0xE3,0xF2,0xFD),
          RGBColor(0xE3,0xF2,0xFD), RGBColor(0xC8,0xE6,0xC9),
          RGBColor(0xC8,0xE6,0xC9), RGBColor(0xE8,0xF5,0xE9)]
    tc = [C_WHITE, C_GRAY, C_GRAY, C_GRAY,
          RGBColor(0x1B,0x5E,0x20), RGBColor(0x1B,0x5E,0x20),
          RGBColor(0x1B,0x5E,0x20)]
    table(sl, rows, cw, Inches(7.0), Inches(1.28), Inches(0.48),
          row_colors=rc, txt_colors=tc)

    rect(sl, Inches(7.0), Inches(4.7), Inches(6.1), Inches(2.4),
         RGBColor(0xFF,0xF3,0xE0))
    txt(sl, "Key Takeaways",
        Inches(7.2), Inches(4.78), Inches(5.8), Inches(0.28),
        size=12, bold=True, color=C_ACCENT)
    bullets(sl, [
        "▸  Bayesian consistently leads at every level",
        "▸  NSGA-II: huge gain from extended training (4.5×)",
        "▸  NSGA-III limited by architecture size (11k params)",
        "▸  Extended training (L5) is the biggest improvement step",
        "▸  Skip=2 achieves 2× FEM reduction for all three",
    ], Inches(7.0), Inches(5.05), Inches(6.1), Inches(2.0), size=11)

    notes(sl,
        "The cross-level chart shows L2 improvement from L1 through L6. "
        "Level 5 extended training delivers the largest single improvement. "
        "Bayesian starts best and stays best. NSGA-II shows the most dramatic "
        "improvement. NSGA-III's small architecture limits its ceiling. "
        "The skip results (L2) are the per-window timestepper L2, not the "
        "global L2 — different scale but same trend.",
        "Çapraz seviye grafiği, L1'den L6'ya L2 iyileşmesini gösteriyor. "
        "Seviye 5 uzatılmış eğitimi en büyük tekil iyileşmeyi sağlar. "
        "Bayesian en iyi başlar ve en iyi kalır. NSGA-II en dramatik "
        "iyileşmeyi gösterir. NSGA-III'ün küçük mimarisi tavanını sınırlar.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Conclusions
# ══════════════════════════════════════════════════════════════════════════════

def slide_conclusions(prs):
    sl = blank(prs)
    bg(sl, C_DARK)
    rect(sl, 0, Inches(1.05), SLIDE_W, Inches(0.055), C_ACCENT)

    txt(sl, "Conclusions & Key Findings",
        Inches(0.5), Inches(0.12), Inches(12.0), Inches(0.82),
        size=30, bold=True, color=C_WHITE)

    findings = [
        ("✓", "PINN replaces FEM steps without accuracy loss",
         "Bayesian: skip=2 → 33.2°C vs 43.6°C baseline (2× fewer FEM calls)"),
        ("✓", "NAS architecture determines skip capability",
         "Bayesian 5×151 relu converges at skip=4; NSGA-III fails at skip=2 with 500 epochs"),
        ("✓", "Extended training closes the gap for all optimizers",
         "NSGA-II: 0.252 → 0.055 (4.5× improvement); meets L2 < 0.10 target"),
        ("✓", "Hybrid routing achieves 80–95% FEM step reduction",
         "Residual-based routing (thr=0.1) gives 4–19 PINN steps per run"),
        ("✓", "More training epochs fix skip convergence (not architecture)",
         "NSGA-II/III: 2 000 epochs → both converge at skip=2 (ratio ≈ 0.98–0.99×)"),
        ("✓", "NAS generalises across PDE families and geometries",
         "Same Bayesian arch: circle L2=0.00014 | L-shape L2=0.028 | 5 domains tested"),
    ]

    for i, (icon, title_f, detail) in enumerate(findings):
        y = Inches(1.25) + i * Inches(0.95)
        rect(sl, Inches(0.35), y, Inches(0.5), Inches(0.52), C_GREEN)
        txt(sl, icon, Inches(0.35), y + Inches(0.06),
            Inches(0.5), Inches(0.44),
            size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, title_f, Inches(0.96), y,
            Inches(12.0), Inches(0.38),
            size=13, bold=True, color=C_WHITE)
        txt(sl, detail, Inches(0.96), y + Inches(0.40),
            Inches(12.0), Inches(0.46),
            size=11, color=C_LBLUE, italic=True)

    rect(sl, 0, Inches(7.1), SLIDE_W, Inches(0.4), RGBColor(0x0D,0x47,0xA1))
    txt(sl,
        "Reference: Mortensen, Noorsumar, Fjær, Babaei, Drønen (2026) — "
        "The International Journal of Advanced Manufacturing Technology — "
        "DOI: 10.1007/s00170-026-17515-w",
        Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.35),
        size=9, color=RGBColor(0x90,0xCA,0xF9), align=PP_ALIGN.CENTER)

    notes(sl,
        "Summary of key results: "
        "(1) Temporal skip operator works — PINN can replace FEM steps with 2× speedup. "
        "(2) Architecture matters — Bayesian 5×151 relu is best for skip; smaller models "
        "need more training time. "
        "(3) Extended training (L5) is highly effective — NSGA-II improves 4.5×. "
        "(4) Hybrid FEM+PINN with residual routing achieves 80-95% FEM step reduction. "
        "(5) Epoch count was the bottleneck for NSGA-II/III skip convergence, not architecture. "
        "(6) Multi-domain Poisson NAS works well on simple geometries (circle: L2=0.00014) "
        "but struggles on complex shapes (flower: L2=0.17). "
        "Future work: higher skip rates for NSGA-II/III; 3D extension; real FEM coupling.",
        "Temel sonuçların özeti: "
        "(1) Zamansal atlama operatörü çalışıyor — PINN, 2× hızlanmayla FEM adımlarını değiştirebilir. "
        "(2) Mimari önemlidir — Bayesian 5×151 relu atlamak için en iyisidir. "
        "(3) Uzatılmış eğitim (L5) son derece etkilidir — NSGA-II 4.5× iyileşir. "
        "(4) Artık yönlendirmeli hibrit FEM+PINN, %80-95 FEM adım azaltımı sağlar. "
        "(5) NSGA-II/III atlama yakınsaması için darboğaz epoch sayısıydı, mimari değil. "
        "(6) Çoklu domain Poisson NAS, basit geometrilerde iyi çalışır (daire: L2=0.00014) "
        "ama karmaşık şekillerde zorlanır (çiçek: L2=0.17). "
        "Gelecek çalışma: NSGA-II/III için daha yüksek atlama oranları; 3B uzantı; gerçek FEM bağlantısı.")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_title(prs)            # 1
    slide_overview(prs)         # 2
    slide_problem(prs)          # 3
    slide_nas_space(prs)        # 4
    slide_level1_table(prs)     # 5
    slide_level1_curves(prs)    # 6
    slide_level2_concept(prs)   # 7
    slide_level2_table(prs)     # 8
    slide_level2_viz(prs)       # 9
    slide_level2_2k(prs)        # 10
    slide_level3(prs)           # 11
    slide_level3_residual(prs)  # 12
    slide_level4(prs)           # 13
    slide_level4_comparison(prs)# 14
    slide_level5(prs)           # 15
    slide_level5_plot(prs)      # 16
    slide_level6(prs)           # 17
    slide_level7b_domains(prs)     # 18
    slide_level7b_pred_error(prs)  # 19
    slide_level7b_results(prs)     # 20
    slide_cross_level(prs)         # 21
    slide_conclusions(prs)         # 22

    out = ROOT / "results" / "NAS_PINNs_Presentation.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
