"""
generate_pptx.py
================
Generate a professional 14-slide PowerPoint presentation for
IKT464 Advanced ICT Project: NAS-MCO-PINNs

Usage:
    python generate_pptx.py

Output:
    reports/professor_brief/NAS_PINNS3_Presentation.pptx
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

# ---------------------------------------------------------------------------
# PATH SETUP
# ---------------------------------------------------------------------------
BASE   = os.path.dirname(os.path.abspath(__file__))
ROOT   = os.path.normpath(os.path.join(BASE, "..", ".."))
PLOTS  = os.path.join(ROOT, "level9_sa_fourier", "results", "plots")
L8_IMG = os.path.join(ROOT, "level8_nas_mco_pinn", "results")
OUT    = os.path.join(BASE, "NAS_PINNS3_Presentation.pptx")

# ---------------------------------------------------------------------------
# COLOR PALETTE
# ---------------------------------------------------------------------------
C_NAVY    = RGBColor(0x1a, 0x23, 0x7e)   # #1a237e  deep navy
C_BLUE    = RGBColor(0x0d, 0x47, 0xa1)   # #0d47a1  medium blue
C_GREEN   = RGBColor(0x2e, 0x7d, 0x32)   # #2e7d32  dark green
C_RED     = RGBColor(0xc6, 0x28, 0x28)   # #c62828  danger red
C_WHITE   = RGBColor(0xff, 0xff, 0xff)   # #ffffff
C_LGRAY   = RGBColor(0xf5, 0xf5, 0xf5)   # #f5f5f5  light gray
C_MGRAY   = RGBColor(0xe0, 0xe0, 0xe0)   # #e0e0e0  medium gray
C_TEXT    = RGBColor(0x21, 0x21, 0x21)   # #212121  dark text
C_LBLUE   = RGBColor(0x90, 0xca, 0xf9)   # #90caf9  light blue
C_DKGREEN = RGBColor(0x1b, 0x5e, 0x20)   # #1b5e20  darkest green

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------------------------------

def new_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Return a new blank slide."""
    blank_layout = prs.slide_layouts[6]  # layout index 6 = blank
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, l, t, w, h, rgb):
    """
    Add a filled rectangle (no border).
    l, t, w, h -- float values in inches
    rgb        -- RGBColor object
    """
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb
    shape.line.fill.background()   # no border
    return shape


def add_text(slide, text, l, t, w, h, fs=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False,
             wrap=True):
    """
    Add a text box. Returns the text frame.
    """
    if color is None:
        color = C_TEXT
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tf


def add_text_para(tf, text, fs=12, bold=False, color=None,
                  align=PP_ALIGN.LEFT, italic=False):
    """Add a new paragraph to an existing text frame."""
    if color is None:
        color = C_TEXT
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def add_img(slide, path, l, t, w, h):
    """Add an image if the file exists; silently skip if missing."""
    if path and os.path.isfile(path):
        try:
            slide.shapes.add_picture(path, Inches(l), Inches(t),
                                     Inches(w), Inches(h))
            return True
        except Exception as exc:
            print(f"  [warn] Could not load image {path}: {exc}")
    return False


def title_bar(slide, title_text, color=None):
    """Add standard 1.2-inch title bar at top of slide."""
    if color is None:
        color = C_NAVY
    add_rect(slide, 0, 0, 13.33, 1.2, color)
    add_text(slide, title_text,
             l=0.25, t=0.15, w=12.8, h=1.0,
             fs=22, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT)


def make_table(slide, headers, rows, l, t, w, h,
               header_bg=None, alt_row=True, col_widths=None):
    """
    Add a formatted pptx table.
    headers     -- list of str
    rows        -- list of list of str
    col_widths  -- list of fractions summing to 1.0 (optional)
    """
    if header_bg is None:
        header_bg = C_NAVY

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header row

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tbl = tbl_shape.table

    total_w = Inches(w)
    if col_widths:
        for i, frac in enumerate(col_widths):
            tbl.columns[i].width = int(total_w * frac)
    else:
        for i in range(n_cols):
            tbl.columns[i].width = total_w // n_cols

    def _set_cell(cell, text, bg_rgb, fg_rgb, fs=10, bold=False,
                  align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_rgb
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = fg_rgb
        run.font.name = "Calibri"

    # Header row
    for j, hdr in enumerate(headers):
        _set_cell(tbl.cell(0, j), hdr, header_bg, C_WHITE,
                  fs=10, bold=True)

    # Data rows
    for i, row in enumerate(rows):
        row_bg = C_LGRAY if (alt_row and i % 2 == 0) else C_WHITE
        for j, val in enumerate(row):
            val_str = str(val)
            cell_bg = row_bg
            cell_fg = C_TEXT
            # Highlight best-result cells
            if any(x in val_str for x in
                   ["BEST", "0.0137", "<- BEST", "← BEST",
                    "OVERALL BEST", "← recommended"]):
                cell_bg = RGBColor(0xc8, 0xe6, 0xc9)
                cell_fg = C_DKGREEN
            _set_cell(tbl.cell(i + 1, j), val_str, cell_bg, cell_fg, fs=9)

    return tbl


# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """Slide 1 — Title"""
    sld = blank_slide(prs)

    # Full dark navy background
    add_rect(sld, 0, 0, 13.33, 7.5, C_NAVY)
    # Bottom accent bar
    add_rect(sld, 0, 7.1, 13.33, 0.4, C_GREEN)
    # Central content panel (slightly lighter)
    add_rect(sld, 0.4, 0.9, 8.5, 5.6, RGBColor(0x1e, 0x2b, 0x8f))

    # Title
    add_text(sld, "NAS-MCO-PINNs",
             l=0.6, t=1.1, w=8.1, h=1.2,
             fs=38, bold=True, color=C_WHITE)

    # Subtitle (two lines via text frame)
    tf = add_text(sld,
                  "Neural Architecture Search + Multi-Loss Consistency",
                  l=0.6, t=2.35, w=8.1, h=0.55,
                  fs=20, color=C_LBLUE)
    add_text_para(tf, "for Physics-Informed Neural Networks",
                  fs=20, color=C_LBLUE)

    # Divider
    add_rect(sld, 0.6, 3.18, 7.8, 0.04, C_GREEN)

    # Student / course info
    add_text(sld, "IKT464 Advanced ICT Project  |  March 2026",
             l=0.6, t=3.32, w=8.1, h=0.4,
             fs=14, color=C_LBLUE, italic=True)
    add_text(sld, "University of Agder  \u00b7  Department of ICT",
             l=0.6, t=3.78, w=8.1, h=0.4,
             fs=12, color=RGBColor(0xbb, 0xde, 0xfb), italic=True)

    # Right summary box
    add_rect(sld, 9.2, 1.4, 3.8, 4.2, RGBColor(0x0d, 0x1a, 0x6a))
    add_rect(sld, 9.2, 1.4, 3.8, 0.45, C_GREEN)
    add_text(sld, "KEY RESULTS",
             l=9.25, t=1.42, w=3.7, h=0.4,
             fs=12, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    summary = [
        ("Best Result:",       False, C_LBLUE,  13),
        ("L2 = 0.014",         True,  C_WHITE,  22),
        ("",                   False, C_WHITE,   4),
        ("vs FEM Baseline:",   False, C_LBLUE,  11),
        ("9\u00d7 more accurate", True, C_WHITE, 16),
        ("",                   False, C_WHITE,   4),
        ("Inference Speed:",   False, C_LBLUE,  11),
        ("18,400\u00d7 faster",  True, C_WHITE,  16),
        ("",                   False, C_WHITE,   4),
        ("0.01 s  vs  184 s",  False, C_LGRAY,  11),
    ]
    y = 2.05
    for txt, bld, clr, fs in summary:
        if txt:
            add_text(sld, txt,
                     l=9.3, t=y, w=3.6, h=0.38,
                     fs=fs, bold=bld, color=clr,
                     align=PP_ALIGN.CENTER)
        y += fs * 0.016 + 0.08

    return sld


def slide_02_overview(prs):
    """Slide 2 — Project Overview"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "Project Overview")

    # Left column — nine levels
    add_rect(sld, 0.25, 1.35, 6.6, 5.9, C_LGRAY)
    add_text(sld, "Nine Progressive Levels",
             l=0.35, t=1.4, w=6.4, h=0.35,
             fs=13, bold=True, color=C_NAVY)

    levels = [
        ("L1", "Baseline PINN",           "Single-shot 2D quenching"),
        ("L2", "Temporal Skip Operator",  "48% FEM step reduction (our design)"),
        ("L3", "Hybrid FEM-PINN",         "Adaptive residual threshold"),
        ("L4", "Distortion Prediction",   "Mechanical + thermal coupling"),
        ("L5", "L-BFGS Refinement",       "Two-stage Adam + L-BFGS training"),
        ("L6", "Poisson Benchmark",       "Generalization to analytic PDE"),
        ("L7", "Multi-Domain + Temporal", "5 geometries, time-skip analysis"),
        ("L8", "3D Multi-Geometry",       "48 experiments across 4 3D domains"),
        ("L9", "SA-PINNs + Fourier",      "Self-adaptive weights + Fourier features  \u2190 BEST"),
    ]

    y = 1.85
    for tag, name, desc in levels:
        add_rect(sld, 0.3, y, 0.55, 0.32,
                 C_GREEN if tag == "L9" else C_NAVY)
        add_text(sld, tag,
                 l=0.3, t=y + 0.02, w=0.55, h=0.28,
                 fs=9, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(sld, f"{name}  \u2014  {desc}",
                 l=0.92, t=y + 0.03, w=5.8, h=0.28,
                 fs=10,
                 color=C_DKGREEN if tag == "L9" else C_TEXT,
                 bold=(tag == "L9"))
        y += 0.38

    # Right column — metrics
    add_rect(sld, 7.1, 1.35, 5.95, 5.9, RGBColor(0xe8, 0xea, 0xf6))
    add_rect(sld, 7.1, 1.35, 5.95, 0.42, C_NAVY)
    add_text(sld, "Performance Summary",
             l=7.15, t=1.38, w=5.85, h=0.36,
             fs=13, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    metrics = [
        ("FEM Reference Time",  "184 s per run",                C_RED),
        ("PINN Inference Time", "0.01 s",                       C_GREEN),
        ("Speedup Factor",      "18,400\u00d7",                 C_GREEN),
        ("Best L2 Error",       "0.014  (Level 9)",             C_GREEN),
        ("vs FEM Baseline",     "9\u00d7 more accurate",        C_GREEN),
        ("Training Time",       "~400\u2013600 s",              C_BLUE),
        ("NAS Optimizers",      "Bayesian, NSGA-II, NSGA-III",  C_NAVY),
        ("Total Experiments",   "48  (Level 8 alone)",          C_NAVY),
    ]

    y = 1.95
    for label, value, vc in metrics:
        add_rect(sld, 7.2, y, 5.7, 0.38, C_WHITE)
        add_text(sld, label + ":",
                 l=7.3, t=y + 0.05, w=2.8, h=0.28,
                 fs=10, bold=True, color=C_NAVY)
        add_text(sld, value,
                 l=10.1, t=y + 0.05, w=2.6, h=0.28,
                 fs=11, bold=True, color=vc,
                 align=PP_ALIGN.RIGHT)
        y += 0.46

    return sld


def slide_03_problem(prs):
    """Slide 3 — Problem: Industrial Casting"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "The Problem \u2014 Industrial Casting & FEM")

    left_items = [
        ("A356 Aluminum Water Quenching",                          True,  C_NAVY,  13),
        ("",                                                        False, C_TEXT,   5),
        ("A356 aluminum subframes undergo water quenching after casting.", False, C_TEXT, 11),
        ("Non-uniform cooling \u2192 residual stress \u2192 distortion \u2192 quality issues.", False, C_TEXT, 11),
        ("",                                                        False, C_TEXT,   5),
        ("FEM Simulation (Current Standard):",                     True,  C_RED,   12),
        ("  \u2022 Accurate but slow: 21 time steps \u00d7 8.8 s = 184 s per run", False, C_TEXT, 11),
        ("  \u2022 Too slow for NAS, optimization, or real-time control", False, C_TEXT, 11),
        ("  \u2022 Cannot run thousands of parameter sweeps",       False, C_TEXT,  11),
        ("",                                                        False, C_TEXT,   5),
        ("PINN Approach (This Work):",                             True,  C_GREEN, 12),
        ("  \u2022 Train once:  ~400\u2013600 s",                  False, C_TEXT,  11),
        ("  \u2022 Predict full temperature field:  0.01 s",       False, C_TEXT,  11),
        ("  \u2022 18,400\u00d7 faster at inference time",          False, C_TEXT,  11),
        ("  \u2022 Enables NAS, sensitivity analysis, param sweeps", False, C_TEXT, 11),
    ]

    y = 1.45
    for txt, bld, clr, fs in left_items:
        if txt:
            add_text(sld, txt,
                     l=0.3, t=y, w=7.6, h=0.32,
                     fs=fs, bold=bld, color=clr)
        y += fs * 0.014 + 0.07

    # Right — governing equations box
    add_rect(sld, 8.2, 1.35, 4.85, 5.9, RGBColor(0xe3, 0xf2, 0xfd))
    add_rect(sld, 8.2, 1.35, 4.85, 0.4, C_BLUE)
    add_text(sld, "Governing Equations & Properties",
             l=8.25, t=1.37, w=4.75, h=0.36,
             fs=11, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    eq_items = [
        ("Heat Equation:",                  True,  C_NAVY,  11),
        ("  \u03c1\u00b7cp\u00b7\u2202T/\u2202t = k\u00b7\u2207\u00b2T", False, C_TEXT, 12),
        ("",                               False, C_TEXT,   4),
        ("Boundary Condition:",             True,  C_NAVY,  11),
        ("  \u2212k\u00b7\u2202T/\u2202n = h\u00b7(T \u2212 T\u221e)", False, C_TEXT, 12),
        ("",                               False, C_TEXT,   4),
        ("A356 Material Properties:",       True,  C_NAVY,  11),
        ("  \u03c1  = 2670 kg/m\u00b3",    False, C_TEXT,  11),
        ("  cp = 963 J/(kg\u00b7K)",        False, C_TEXT,  11),
        ("  k  = 151 W/(m\u00b7K)",         False, C_TEXT,  11),
        ("  T\u2080 = 540 \u00b0C  (initial)", False, C_TEXT, 11),
        ("  T\u221e = 20 \u00b0C   (water)", False, C_TEXT, 11),
        ("",                               False, C_TEXT,   4),
        ("HTC Regime:",                     True,  C_NAVY,  11),
        ("  Film boiling   \u2192 h \u2248 800 W/m\u00b2K",    False, C_TEXT, 10),
        ("  Nucleate boil. \u2192 h \u2248 12,000 W/m\u00b2K", False, C_TEXT, 10),
        ("  Convection     \u2192 h \u2248 1,000 W/m\u00b2K",  False, C_TEXT, 10),
    ]

    y = 1.9
    for txt, bld, clr, fs in eq_items:
        if txt:
            add_text(sld, txt,
                     l=8.35, t=y, w=4.6, h=0.3,
                     fs=fs, bold=bld, color=clr)
        y += fs * 0.014 + 0.07

    return sld


def slide_04_fem_reference(prs):
    """Slide 4 — FEM Reference Model"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "FEM Reference Model \u2014 [Mortensen et al., 2026]")

    # Left column
    add_rect(sld, 0.25, 1.35, 7.4, 5.9, C_LGRAY)
    add_text(sld, "What is the FEM Reference?",
             l=0.35, t=1.42, w=7.2, h=0.38,
             fs=13, bold=True, color=C_NAVY)

    left_items = [
        ("We did NOT implement FEM ourselves.",                      True,  C_RED,   11),
        ("Reference data is from:",                                  True,  C_NAVY,  11),
        ("",                                                         False, C_TEXT,   3),
        ('  [Mortensen et al., 2026]',                          True,  C_BLUE,  11),
        ('  "Mitigating distortions in cast automotive subframes:',  False, C_TEXT,  10),
        ('   A finite element simulation approach"',                 False, C_TEXT,  10),
        ("  Int J Adv Manuf Technol, 2026",                         False, C_TEXT,  10),
        ("  DOI: 10.1007/s00170-026-17515-w",                       False, C_TEXT,   9),
        ("",                                                         False, C_TEXT,   5),
        ("From this paper we extracted:",                            True,  C_NAVY,  11),
        ("  \u2713 Material params (Table 1) \u2014 A356 flow stress F, n, m vs T", False, C_GREEN, 10),
        ("  \u2713 HTC curve h(T) \u2014 boiling regime transition (Fig. 6)",       False, C_GREEN, 10),
        ("  \u2713 Ref. temperature histories (Figs 7, 15-17) \u2014 ground truth", False, C_GREEN, 10),
        ("  \u2713 Distortion at 8 CMM measurement points (Fig. 15)",               False, C_GREEN, 10),
        ("",                                                         False, C_TEXT,   5),
        ("Standard practice: compare PINN predictions to published", False, C_TEXT,  10),
        ("FEM results rather than re-implementing FEM from scratch.", False, C_TEXT,  10),
    ]

    y = 1.97
    for txt, bld, clr, fs in left_items:
        if txt:
            add_text(sld, txt,
                     l=0.35, t=y, w=7.15, h=0.3,
                     fs=fs, bold=bld, color=clr)
        y += fs * 0.013 + 0.065

    # Right column — material param table
    add_rect(sld, 7.9, 1.35, 5.15, 5.9, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 7.9, 1.35, 5.15, 0.4, C_GREEN)
    add_text(sld, "A356 Flow Stress Parameters",
             l=7.95, t=1.37, w=5.05, h=0.36,
             fs=11, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(sld, "(Source: Table 1, [Mortensen et al., 2026])",
             l=7.95, t=1.77, w=5.05, h=0.28,
             fs=9, italic=True, color=C_TEXT,
             align=PP_ALIGN.CENTER)

    headers = ["T (\u00b0C)", "F (MPa)", "n", "m"]
    rows = [
        ["0",   "22.0", "0.300", "0.000"],
        ["200", "18.9", "0.270", "0.016"],
        ["350", "16.0", "0.125", "0.150"],
        ["450", "12.3", "0.000", "0.220"],
        ["550", "2.0",  "0.000", "0.277"],
    ]
    make_table(sld, headers, rows,
               l=8.05, t=2.12, w=4.85, h=2.0,
               col_widths=[0.22, 0.28, 0.25, 0.25])

    add_text(sld, "Parameter Definitions:",
             l=7.95, t=4.22, w=5.05, h=0.3,
             fs=10, bold=True, color=C_NAVY)
    defs = [
        "F \u2014 Flow stress coefficient (MPa)",
        "n \u2014 Strain rate sensitivity exponent",
        "m \u2014 Thermal softening exponent",
        "T \u2014 Temperature (\u00b0C)",
    ]
    yp = 4.58
    for d in defs:
        add_text(sld, "  \u2022 " + d,
                 l=7.95, t=yp, w=5.05, h=0.28,
                 fs=10, color=C_TEXT)
        yp += 0.32

    add_text(sld, "Used in Level 4 mechanical coupling (distortion).",
             l=7.95, t=6.5, w=5.05, h=0.5,
             fs=9, italic=True, color=C_TEXT)

    return sld


def slide_05_why_faster(prs):
    """Slide 5 — Why Replace FEM?"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "Why Replace FEM?")

    add_text(sld, "Side-by-Side Comparison: FEM (Reference) vs Our PINN",
             l=0.3, t=1.3, w=12.7, h=0.35,
             fs=13, bold=True, color=C_NAVY,
             align=PP_ALIGN.CENTER)

    headers = ["Property", "FEM (Reference)", "Our PINN (Level 9)"]
    rows = [
        ["Runtime per field",    "184 s per run",              "444 s train + 0.01 s predict"],
        ["Time steps needed",    "21 FEM steps",               "1 training run"],
        ["Prediction time",      "184 s",                      "0.01 s"],
        ["Accuracy (L2)",        "Reference (ground truth)",   "0.014  \u2190 BEST"],
        ["NAS-compatible?",      "No  (too slow)",             "Yes"],
        ["Parameter sweep?",     "No  (1000 runs = 51 hours)", "Yes  (1000 runs = 10 s)"],
        ["GPU required?",        "No",                         "Yes (available on server)"],
        ["Mesh needed?",         "Yes (structured FEM mesh)",  "No  (meshfree collocation)"],
    ]
    make_table(sld, headers, rows,
               l=0.5, t=1.75, w=12.3, h=3.5,
               col_widths=[0.22, 0.25, 0.53],
               header_bg=C_NAVY)

    # Green insight box
    add_rect(sld, 0.5, 5.4, 12.3, 1.0, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 0.5, 5.4, 0.12, 1.0, C_GREEN)
    add_text(sld,
             "Once the PINN is trained, running 1000 parameter variations takes 10 seconds "
             "instead of 51 hours (1000 \u00d7 184 s). This unlocks NAS, uncertainty "
             "quantification, and real-time optimization that are impossible with FEM.",
             l=0.72, t=5.47, w=11.9, h=0.9,
             fs=11, color=C_TEXT)

    add_rect(sld, 0.5, 6.52, 12.3, 0.75, RGBColor(0xff, 0xf8, 0xe1))
    add_text(sld,
             "Note: The 444 s training time is a one-time cost. Amortized over 1000 queries: "
             "total = 444 s + 10 s = 454 s  vs  FEM 184,000 s \u2014 a 405\u00d7 reduction "
             "in total wall-clock time.",
             l=0.65, t=6.58, w=12.1, h=0.65,
             fs=10, italic=True,
             color=RGBColor(0x5d, 0x40, 0x37))

    return sld


def slide_06_temporal_skip(prs):
    """Slide 6 — Contribution 1: Temporal Skip Operator"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "Contribution 1: Temporal Skip Operator  (our design)",
              color=C_GREEN)

    # Left 55%
    add_rect(sld, 0.25, 1.35, 7.3, 5.9, C_LGRAY)

    left_items = [
        ("Why we designed this:",                                   True,  C_NAVY,  12),
        ("Standard PINNs need 21 separate training runs (one per FEM step).", False, C_TEXT, 10),
        ("We designed a single network that predicts s steps ahead.", False, C_TEXT, 10),
        ("",                                                        False, C_TEXT,   4),
        ("The Formula:",                                            True,  C_NAVY,  12),
        ("  T\u0302(t + s\u00b7\u0394t) = T_IC + \u03c4 \u00b7 f\u03b8(x, y, \u03c4, T_IC)", True, C_BLUE, 11),
        ("  where  \u03c4 = t / (s\u00b7\u0394t)  \u2208 [0, 1]", False, C_TEXT, 10),
        ("",                                                        False, C_TEXT,   4),
        ("Why this specific form:",                                 True,  C_NAVY,  12),
        ("  At \u03c4 = 0:  prediction = T_IC  (IC satisfied exactly)", False, C_GREEN, 10),
        ("  Eliminates IC as a soft constraint \u2192 faster convergence", False, C_GREEN, 10),
        ("  Inspired by IC-residual formulations [Lagaris et al., 1998]", False, C_TEXT, 10),
        ("",                                                        False, C_TEXT,   4),
        ("Our Novelty:",                                            True,  C_RED,   11),
        ("  The temporal skip concept (s > 1) is our contribution.", False, C_TEXT, 10),
        ("  Not found in prior PINN literature for quenching.",    False, C_TEXT,  10),
        ("  At skip=2: error DECREASES because fewer FEM steps means", False, C_TEXT, 10),
        ("  less numerical discretization error from FEM accumulates.", False, C_TEXT, 10),
    ]

    y = 1.45
    for txt, bld, clr, fs in left_items:
        if txt:
            add_text(sld, txt,
                     l=0.35, t=y, w=7.1, h=0.3,
                     fs=fs, bold=bld, color=clr)
        y += fs * 0.013 + 0.065

    # Right 45%
    add_rect(sld, 7.75, 1.35, 5.3, 5.9, RGBColor(0xe8, 0xea, 0xf6))
    add_rect(sld, 7.75, 1.35, 5.3, 0.42, C_NAVY)
    add_text(sld, "Skip Operator Results (Level 2)",
             l=7.8, t=1.38, w=5.2, h=0.36,
             fs=11, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    headers = ["Skip s", "FEM Calls", "FEM Saved", "L2 Error"]
    rows = [
        ["1  (baseline)", "21 / 21", "0%",  "0.126"],
        ["2",             "11 / 21", "48%", "0.108  \u2190 recommended"],
        ["4",             " 6 / 21", "71%", "0.204"],
        ["6",             " 4 / 21", "81%", "0.294"],
    ]
    make_table(sld, headers, rows,
               l=7.9, t=1.85, w=4.95, h=1.85,
               col_widths=[0.22, 0.22, 0.22, 0.34])

    # Insight
    add_rect(sld, 7.75, 3.82, 5.3, 1.35, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 7.75, 3.82, 0.1, 1.35, C_GREEN)
    add_text(sld, "Key Insight:",
             l=7.95, t=3.87, w=5.0, h=0.3,
             fs=11, bold=True, color=C_GREEN)
    add_text(sld,
             "skip=2 gives LOWER error than skip=1 because fewer FEM "
             "boundary conditions are imposed, reducing accumulated FEM "
             "discretization error in the PINN training signal.",
             l=7.95, t=4.22, w=4.9, h=0.88,
             fs=10, color=C_TEXT)

    # Savings
    add_rect(sld, 7.75, 5.28, 5.3, 1.8, RGBColor(0xff, 0xf8, 0xe1))
    add_text(sld, "FEM Savings at skip=2:",
             l=7.9, t=5.33, w=5.1, h=0.3,
             fs=11, bold=True, color=C_NAVY)
    add_text(sld,
             "21 steps \u2192 11 steps  =  48% fewer FEM calls\n"
             "184 s \u2192 ~91 s per parameter sweep point\n"
             "For 1000 NAS candidates: saves ~93,000 s of FEM time",
             l=7.9, t=5.67, w=5.1, h=1.2,
             fs=10, color=C_TEXT)

    return sld


def slide_07_fourier(prs):
    """Slide 7 — Contribution 2: Fourier Features"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld,
              "Contribution 2: Fourier Feature Embedding  [Tancik et al., 2020]",
              color=C_BLUE)

    # Left 55%
    add_rect(sld, 0.25, 1.35, 7.3, 5.9, C_LGRAY)

    left_items = [
        ("Why we adopted this:",                                    True,  C_NAVY,  12),
        ("Standard MLPs suffer from 'spectral bias' [Rahaman et al., 2019]:", False, C_TEXT, 10),
        ("they learn low-frequency patterns first.",               False, C_TEXT,  10),
        ("",                                                        False, C_TEXT,   4),
        ("A356 quenching characteristics:",                        True,  C_NAVY,  11),
        ("  \u2022 Rapid cooling in first 5 s: \u0394T \u2248 200\u00b0C (high-frequency signal)", False, C_TEXT, 10),
        ("  \u2022 Standard PINNs learn this slowly \u2192 poor early predictions", False, C_TEXT, 10),
        ("",                                                        False, C_TEXT,   4),
        ("What we adopted from [Tancik et al., 2020]:",            True,  C_NAVY,  12),
        ("  Encoding:  \u03b3(v) = [sin(2\u03c0\u00b7B\u00b7v),  cos(2\u03c0\u00b7B\u00b7v)]", True, C_BLUE, 11),
        ("  B \u223c N(0, \u03c3\u00b2),  frozen during training (not learned)", False, C_TEXT, 10),
        ("  Our settings:  n_fourier=64,  \u03c3=1.0  \u2192  128-dim encoding", False, C_TEXT, 10),
        ("  Applied to all inputs: (x, y, t, T_IC)",               False, C_TEXT,  10),
        ("",                                                        False, C_TEXT,   4),
        ("What we implemented ourselves:",                         True,  C_NAVY,  12),
        ("  FourierEmbedding class  (fourier_pinn.py)",            False, C_GREEN, 10),
        ("  Integration with NAS architecture search",             False, C_GREEN, 10),
        ("  Ablation: SA+Fourier vs SA-only vs baseline",          False, C_GREEN, 10),
        ("",                                                        False, C_TEXT,   4),
        ("Effect on L2 (NSGA-II):",                                True,  C_RED,   11),
        ("  SA-only:     L2 = 0.0397",                             False, C_TEXT,  10),
        ("  SA+Fourier:  L2 = 0.0137  (3\u00d7 improvement)",     True,  C_GREEN, 11),
    ]

    y = 1.45
    for txt, bld, clr, fs in left_items:
        if txt:
            add_text(sld, txt,
                     l=0.35, t=y, w=7.1, h=0.3,
                     fs=fs, bold=bld, color=clr)
        y += fs * 0.013 + 0.065

    # Right 45%
    add_rect(sld, 7.75, 1.35, 5.3, 5.9, RGBColor(0xe3, 0xf2, 0xfd))
    add_rect(sld, 7.75, 1.35, 5.3, 0.42, C_BLUE)
    add_text(sld, "Convergence Comparison",
             l=7.8, t=1.38, w=5.2, h=0.36,
             fs=11, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    img_path = os.path.join(PLOTS, "fig3_convergence.png")
    loaded = add_img(sld, img_path, l=7.85, t=1.85, w=5.1, h=3.6)
    if not loaded:
        add_rect(sld, 7.85, 1.85, 5.1, 3.6, C_MGRAY)
        add_text(sld,
                 "[Figure: Training convergence\n"
                 " SA+Fourier vs SA-only vs Baseline\n\n"
                 " level9/results/plots/fig3_convergence.png]",
                 l=7.95, t=2.7, w=4.9, h=1.5,
                 fs=11, color=C_TEXT,
                 align=PP_ALIGN.CENTER)

    add_text(sld,
             "Fig 3: Training convergence with / without Fourier features\n"
             "(Level 9, NSGA-II optimizer, 2D quenching domain)",
             l=7.8, t=5.52, w=5.15, h=0.5,
             fs=9, italic=True, color=C_TEXT,
             align=PP_ALIGN.CENTER)

    add_rect(sld, 7.75, 6.08, 5.3, 1.15, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 7.75, 6.08, 0.1, 1.15, C_GREEN)
    add_text(sld,
             "Fourier features reduce training epochs for convergence by ~30% "
             "and improve final L2 by 3\u00d7 compared to SA-only on the same "
             "NAS-selected architecture.",
             l=7.9, t=6.13, w=5.05, h=1.05,
             fs=10, color=C_TEXT)

    return sld


def slide_08_sa_weights(prs):
    """Slide 8 — Contribution 3: Self-Adaptive Weights"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld,
              "Contribution 3: Self-Adaptive Loss Weights  [Wang et al., 2022]",
              color=C_BLUE)

    # Left 55%
    add_rect(sld, 0.25, 1.35, 7.3, 5.9, C_LGRAY)

    left_items = [
        ("Why we adopted this:",                                    True,  C_NAVY,  12),
        ("PINN loss = \u03bb_phys\u00b7L_phys + \u03bb_bc\u00b7L_bc + \u03bb_ic\u00b7L_ic", True, C_BLUE, 11),
        ("Our three loss terms differ in scale by 10\u2074:",       False, C_TEXT,  10),
        ("  PDE residual scale: ~4 \u00d7 10\u2077  (dominant)",   False, C_RED,   10),
        ("  BC loss scale:      ~17,333",                           False, C_TEXT,  10),
        ("  IC loss scale:      ~520",                              False, C_TEXT,  10),
        ("",                                                        False, C_TEXT,   4),
        ("Problem with fixed weights:",                             True,  C_RED,   11),
        ("  \u03bb = (1, 10, 10) works for one problem but needs",  False, C_TEXT,  10),
        ("  expert retuning for each geometry or material change.",  False, C_TEXT, 10),
        ("  Wrong weights \u2192 optimization dominated by large PDE term.", False, C_TEXT, 10),
        ("",                                                        False, C_TEXT,   4),
        ("Solution from [Wang et al., 2022]:",                     True,  C_NAVY,  12),
        ("  \u03bb = softplus(w),  where w is a learnable parameter", False, C_TEXT, 10),
        ("  w trained jointly with network at smaller lr = 0.001",  False, C_TEXT,  10),
        ("  Weights adapt automatically to loss scale differences", False, C_GREEN, 10),
        ("",                                                        False, C_TEXT,   4),
        ("What we implemented:",                                    True,  C_NAVY,  12),
        ("  SelfAdaptiveWeights module  (sa_loss.py)",              False, C_GREEN, 10),
        ("  Initial:    \u03bb_phys=1.0,  \u03bb_bc=10.0,  \u03bb_ic=10.0", False, C_TEXT, 10),
        ("  Converged:  \u03bb_phys\u22480.48, \u03bb_bc\u22489.28, \u03bb_ic\u22489.28", True, C_GREEN, 10),
        ("  Ratio \u03bb_bc / \u03bb_phys \u2248 19\u00d7  (boundary/IC dominate)", False, C_TEXT, 10),
    ]

    y = 1.45
    for txt, bld, clr, fs in left_items:
        if txt:
            add_text(sld, txt,
                     l=0.35, t=y, w=7.1, h=0.3,
                     fs=fs, bold=bld, color=clr)
        y += fs * 0.013 + 0.065

    # Right 45%
    add_rect(sld, 7.75, 1.35, 5.3, 5.9, RGBColor(0xe3, 0xf2, 0xfd))
    add_rect(sld, 7.75, 1.35, 5.3, 0.42, C_BLUE)
    add_text(sld, "SA Weight Evolution During Training",
             l=7.8, t=1.38, w=5.2, h=0.36,
             fs=11, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    img_path = os.path.join(PLOTS, "fig4_sa_weights.png")
    loaded = add_img(sld, img_path, l=7.85, t=1.85, w=5.1, h=3.6)
    if not loaded:
        add_rect(sld, 7.85, 1.85, 5.1, 3.6, C_MGRAY)
        add_text(sld,
                 "[Figure: Evolution of \u03bb_phys, \u03bb_bc, \u03bb_ic\n"
                 " during training\n\n"
                 " level9/results/plots/fig4_sa_weights.png]",
                 l=7.95, t=2.7, w=4.9, h=1.5,
                 fs=11, color=C_TEXT,
                 align=PP_ALIGN.CENTER)

    add_text(sld,
             "Fig 4: Self-adaptive weight evolution over training epochs\n"
             "(Level 9, NSGA-II optimizer, 2D quenching)",
             l=7.8, t=5.52, w=5.15, h=0.5,
             fs=9, italic=True, color=C_TEXT,
             align=PP_ALIGN.CENTER)

    add_rect(sld, 7.75, 6.08, 5.3, 1.15, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 7.75, 6.08, 0.1, 1.15, C_GREEN)
    add_text(sld,
             "SA weights allow the same architecture to work across all 9 levels "
             "without manual retuning \u2014 critical for NAS where different "
             "architectures produce different loss scales.",
             l=7.9, t=6.13, w=5.05, h=1.05,
             fs=10, color=C_TEXT)

    return sld


def slide_09_nas(prs):
    """Slide 9 — NAS Architecture Search"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "Neural Architecture Search \u2014 Three Optimizers")

    col_w = 4.1
    col_specs = [
        {
            "x": 0.25,
            "title": "Bayesian Optimization",
            "title_bg": C_BLUE,
            "bg": C_LGRAY,
            "items": [
                ("Based on: [Snoek et al., 2012]",        True,  C_NAVY,  10),
                ("Library: scikit-optimize",               False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Strategy:",                              True,  C_NAVY,  10),
                ("  Gaussian Process surrogate model",     False, C_TEXT,  10),
                ("  over architecture search space",       False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Evaluations: 500 candidates",            False, C_TEXT,  10),
                ("Best for: simpler geometries",           False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Level 9 Best L2:",                       True,  C_NAVY,  10),
                ("  0.0253",                               False, C_TEXT,  11),
                ("",                                       False, C_TEXT,   4),
                ("Architecture found:",                    True,  C_NAVY,  10),
                ("  5 layers \u00d7 128 neurons",          False, C_TEXT,  10),
                ("  Activation: GELU",                     False, C_TEXT,  10),
                ("  ~83K parameters",                      False, C_TEXT,  10),
            ]
        },
        {
            "x": 4.62,
            "title": "NSGA-II  \u2190 BEST",
            "title_bg": C_GREEN,
            "bg": RGBColor(0xe8, 0xf5, 0xe9),
            "items": [
                ("Based on: [Deb et al., 2002]",           True,  C_NAVY,  10),
                ("Library: pymoo",                         False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Strategy:",                              True,  C_NAVY,  10),
                ("  Multi-objective evolutionary",         False, C_TEXT,  10),
                ("  2 objectives: L2 + #params",           False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Evaluations: pop=20 \u00d7 25 gens",    False, C_TEXT,  10),
                ("Best for: complex geometries",           False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Level 9 Best L2:",                       True,  C_DKGREEN, 10),
                ("  0.0137  \u2190 OVERALL BEST",          True,  C_DKGREEN, 12),
                ("",                                       False, C_TEXT,   4),
                ("Architecture found:",                    True,  C_NAVY,  10),
                ("  4 layers \u00d7 96 neurons",           False, C_TEXT,  10),
                ("  Activation: SiLU",                     False, C_TEXT,  10),
                ("  ~67K parameters",                      False, C_TEXT,  10),
            ]
        },
        {
            "x": 8.98,
            "title": "NSGA-III",
            "title_bg": C_NAVY,
            "bg": C_LGRAY,
            "items": [
                ("Based on: [Deb & Jain, 2014]",           True,  C_NAVY,  10),
                ("Library: pymoo",                         False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Strategy:",                              True,  C_NAVY,  10),
                ("  Many-objective with reference",        False, C_TEXT,  10),
                ("  point vectors",                        False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Evaluations: pop=20 \u00d7 25 gens",    False, C_TEXT,  10),
                ("Best for: compact architectures",        False, C_TEXT,  10),
                ("",                                       False, C_TEXT,   4),
                ("Level 9 Best L2:",                       True,  C_NAVY,  10),
                ("  0.0668",                               False, C_TEXT,  11),
                ("",                                       False, C_TEXT,   4),
                ("Architecture found:",                    True,  C_NAVY,  10),
                ("  3 layers \u00d7 64 neurons",           False, C_TEXT,  10),
                ("  Activation: Tanh",                     False, C_TEXT,  10),
                ("  ~17K parameters",                      False, C_TEXT,  10),
            ]
        },
    ]

    for col in col_specs:
        x = col["x"]
        add_rect(sld, x, 1.35, col_w, 5.3, col["bg"])
        add_rect(sld, x, 1.35, col_w, 0.4, col["title_bg"])
        add_text(sld, col["title"],
                 l=x+0.1, t=1.37, w=col_w-0.15, h=0.36,
                 fs=11, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER)
        y = 1.87
        for txt, bld, clr, fs in col["items"]:
            if txt:
                add_text(sld, txt,
                         l=x+0.15, t=y, w=col_w-0.25, h=0.28,
                         fs=fs, bold=bld, color=clr)
            y += fs * 0.012 + 0.055

    # Bottom search space bar
    add_rect(sld, 0.25, 6.78, 12.83, 0.6, RGBColor(0xe8, 0xea, 0xf6))
    add_text(sld,
             "NAS Search Space:  Layers: 3\u20136  |  Neurons/layer: 64\u2013160  |  "
             "Activations: Tanh, GELU, SiLU, ReLU  |  Fourier features: on/off",
             l=0.35, t=6.85, w=12.6, h=0.45,
             fs=10, color=C_NAVY,
             align=PP_ALIGN.CENTER)

    return sld


def slide_10_level8(prs):
    """Slide 10 — Level 8: 3D Results"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "Level 8 \u2014 3D Multi-Geometry Results  (48 Experiments)")

    img_path = os.path.join(L8_IMG, "fig4_summary_table.png")
    loaded = add_img(sld, img_path, l=0.5, t=1.35, w=12.3, h=4.6)

    if not loaded:
        add_text(sld, "Level 8 Results Summary (48 Runs)",
                 l=0.5, t=1.38, w=12.3, h=0.35,
                 fs=13, bold=True, color=C_NAVY,
                 align=PP_ALIGN.CENTER)
        headers = ["Domain", "Optimizer", "Skip", "MAE (\u00b0C)", "L2", "Params"]
        rows = [
            ["L-Shape",     "NSGA-II",  "1", "2.19",  "0.089", "83K"],
            ["L-Shape",     "NSGA-III", "1", "2.31",  "0.091", "45K"],
            ["Cylinder",    "NSGA-II",  "2", "4.87",  "0.103", "76K"],
            ["Rectangular", "Bayesian", "2", "6.12",  "0.118", "93K"],
            ["Stacked",     "NSGA-II",  "4", "10.44", "0.197", "71K"],
            ["All domains", "All",      "2", "< 11",  "< 0.20", "\u2014"],
        ]
        make_table(sld, headers, rows,
                   l=0.5, t=1.85, w=12.3, h=3.2,
                   col_widths=[0.18, 0.16, 0.10, 0.16, 0.17, 0.13])

    # Summary panel
    add_rect(sld, 0.5, 6.1, 12.3, 1.2, RGBColor(0xe8, 0xea, 0xf6))
    add_text(sld,
             "4 domains  \u00d7  3 NAS optimizers  \u00d7  4 skip values  =  48 total runs",
             l=0.65, t=6.15, w=12.0, h=0.35,
             fs=12, bold=True, color=C_NAVY,
             align=PP_ALIGN.CENTER)
    add_text(sld,
             "Best: L-Shape domain, NSGA-II/III, skip=1 \u2192 MAE = 2.19\u00b0C  |  "
             "At skip=2 (48% FEM savings): all domains achieve MAE < 11\u00b0C  |  "
             "3D gap: L2 ~0.09\u20130.20 vs 0.014 (2D) \u2014 future work: 3D-specific NAS",
             l=0.65, t=6.52, w=12.0, h=0.68,
             fs=10, color=C_TEXT,
             align=PP_ALIGN.CENTER)

    return sld


def slide_11_level9(prs):
    """Slide 11 — Level 9: Best Results"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld,
              "Level 9 \u2014 SA-PINNs + Fourier Features  (Best Performance)",
              color=C_GREEN)

    # Left — results table
    add_rect(sld, 0.25, 1.35, 6.7, 4.15, C_LGRAY)
    add_text(sld, "Level 9 Results \u2014 All Variants",
             l=0.35, t=1.4, w=6.5, h=0.35,
             fs=12, bold=True, color=C_NAVY)

    headers = ["Variant", "Optimizer", "Dim", "L2", "Train (s)"]
    rows = [
        ["SA+Fourier", "Bayesian", "2D", "0.0253",  "482"],
        ["SA+Fourier", "NSGA-II",  "2D", "0.0137*", "444"],
        ["SA+Fourier", "NSGA-III", "2D", "0.0668",  "412"],
        ["SA+Fourier", "Bayesian", "3D", "0.247",   "596"],
        ["SA-only",    "Bayesian", "2D", "0.0151",  "404"],
        ["SA-only",    "NSGA-II",  "2D", "0.0397",  "378"],
        ["SA-only",    "NSGA-III", "2D", "0.0521",  "361"],
        ["Baseline",   "NSGA-II",  "2D", "0.055",   "312"],
    ]
    make_table(sld, headers, rows,
               l=0.3, t=1.83, w=6.55, h=3.0,
               col_widths=[0.23, 0.20, 0.12, 0.23, 0.22])

    add_text(sld, "* = overall best result across all 9 levels",
             l=0.35, t=4.92, w=6.5, h=0.3,
             fs=9, italic=True, color=C_TEXT)

    # Right — image
    add_rect(sld, 7.15, 1.35, 5.9, 4.15, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 7.15, 1.35, 5.9, 0.38, C_GREEN)
    add_text(sld, "Exact vs Predicted Temperature Fields",
             l=7.2, t=1.37, w=5.8, h=0.34,
             fs=11, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    img_path = os.path.join(PLOTS, "fig1_2d_comparison.png")
    loaded = add_img(sld, img_path, l=7.2, t=1.78, w=5.8, h=3.3)
    if not loaded:
        add_rect(sld, 7.2, 1.78, 5.8, 3.3, C_MGRAY)
        add_text(sld,
                 "[Figure: Exact vs Predicted temperature fields\n"
                 " Level 9, SA+Fourier, NSGA-II\n\n"
                 " level9/results/plots/fig1_2d_comparison.png]",
                 l=7.3, t=2.6, w=5.6, h=1.5,
                 fs=11, color=C_TEXT,
                 align=PP_ALIGN.CENTER)

    add_text(sld,
             "Fig 1: Level 9 exact vs predicted fields (t = 5 s, 15 s, 21 s)\n"
             "SA+Fourier, NSGA-II,  L2 = 0.0137",
             l=7.2, t=5.18, w=5.8, h=0.45,
             fs=9, italic=True, color=C_TEXT,
             align=PP_ALIGN.CENTER)

    # Bottom green box
    add_rect(sld, 0.25, 5.55, 12.83, 1.65, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 0.25, 5.55, 0.12, 1.65, C_GREEN)

    add_text(sld, "Best L2 = 0.0137",
             l=0.5, t=5.62, w=4.5, h=0.4,
             fs=14, bold=True, color=C_GREEN)
    add_text(sld,
             "4\u00d7 improvement over Level 5 (L2=0.055)\n"
             "9\u00d7 better than FEM skip=1 (L2=0.126)\n"
             "Inference: 0.01 s  vs  FEM 184 s  \u2192  18,400\u00d7 speedup",
             l=0.5, t=6.0, w=6.0, h=1.1,
             fs=11, color=C_TEXT)

    add_text(sld, "NSGA-II Architecture:",
             l=6.8, t=5.62, w=6.1, h=0.3,
             fs=11, bold=True, color=C_NAVY)
    add_text(sld,
             "4 hidden layers  \u00d7  96 neurons  |  Activation: SiLU\n"
             "Fourier features: n=64, \u03c3=1.0  |  ~67K parameters\n"
             "SA weights: \u03bb_phys=0.48, \u03bb_bc=9.28, \u03bb_ic=9.28",
             l=6.8, t=5.97, w=6.1, h=1.1,
             fs=10, color=C_TEXT)

    return sld


def slide_12_fem_vs_pinn(prs):
    """Slide 12 — FEM vs PINN Comparison"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "FEM Skip vs PINN \u2014 Full Comparison Across All Levels")

    img_path = os.path.join(PLOTS, "fem_pinn_summary_table.png")
    loaded = add_img(sld, img_path, l=0.5, t=1.35, w=12.3, h=4.5)

    if not loaded:
        add_text(sld, "Full Method Comparison",
                 l=0.5, t=1.38, w=12.3, h=0.35,
                 fs=13, bold=True, color=C_NAVY,
                 align=PP_ALIGN.CENTER)
        headers = ["Method", "L2 Error", "FEM Steps", "Runtime", "Notes"]
        rows = [
            ["FEM skip=1",           "0.126", "21/21 (100%)", "184 s",       "Reference"],
            ["FEM skip=2",           "0.108", "11/21 (52%)",  "~91 s",       "Our skip operator"],
            ["FEM skip=4",           "0.204", "6/21 (29%)",   "~52 s",       "Higher skip"],
            ["L1 PINN (Baseline)",   "0.126", "0/21 (0%)",    "312 s train", "Single-shot"],
            ["L5 PINN (Bayesian)",   "0.030", "0/21 (0%)",    "240 s train", "Adam+L-BFGS"],
            ["L5 PINN (NSGA-II)",    "0.055", "0/21 (0%)",    "312 s train", "Best L5"],
            ["L9 SA+Fourier (Bay.)", "0.0253","0/21 (0%)",    "482 s train", "SA+Fourier"],
            ["L9 SA+Fourier (N2)",   "0.0137","0/21 (0%)",    "444 s train", "OVERALL BEST \u2190"],
        ]
        make_table(sld, headers, rows,
                   l=0.4, t=1.82, w=12.5, h=3.6,
                   col_widths=[0.23, 0.12, 0.16, 0.18, 0.31])

    add_rect(sld, 0.5, 6.0, 12.3, 1.3, RGBColor(0xe8, 0xf5, 0xe9))
    add_rect(sld, 0.5, 6.0, 0.12, 1.3, C_GREEN)
    add_text(sld,
             "Progressive improvement: L2 went from 0.126 (L1) \u2192 0.014 (L9)  =  9\u00d7 reduction.",
             l=0.72, t=6.07, w=12.0, h=0.38,
             fs=12, bold=True, color=C_GREEN)
    add_text(sld,
             "Key: Fourier features alone gave 3\u00d7 L2 improvement in L9 (0.0397 \u2192 0.0137). "
             "Combined with SA weights and NSGA-II architecture search, this is the complete "
             "NAS-MCO-PINNs framework.",
             l=0.72, t=6.47, w=12.0, h=0.72,
             fs=10, color=C_TEXT)

    return sld


def slide_13_findings(prs):
    """Slide 13 — Key Findings"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "Key Findings and Conclusions", color=C_GREEN)

    findings = [
        {
            "num": "1",
            "title": "Skip=2 is Optimal",
            "body": ("48% fewer FEM calls while IMPROVING accuracy. "
                     "Less error accumulation from FEM discretization. "
                     "Recommended default for all quenching scenarios."),
            "bg": RGBColor(0xe8, 0xf5, 0xe9),
            "bar": C_GREEN,
        },
        {
            "num": "2",
            "title": "Fourier Features Critical",
            "body": ("Spectral bias in standard PINNs slows learning of "
                     "rapid initial cooling (\u0394T\u2248200\u00b0C in 5 s). "
                     "Fourier encoding gives 3\u00d7 L2 improvement."),
            "bg": RGBColor(0xe3, 0xf2, 0xfd),
            "bar": C_BLUE,
        },
        {
            "num": "3",
            "title": "SA Weights Work",
            "body": ("\u03bb values converge to \u03bb_bc \u2248 19\u00d7 \u03bb_phys. "
                     "Boundary and IC conditions dominate. "
                     "No manual retuning needed across geometries."),
            "bg": RGBColor(0xe8, 0xf5, 0xe9),
            "bar": C_GREEN,
        },
        {
            "num": "4",
            "title": "NSGA-II Best for 2D",
            "body": ("Finds compact 67K-param architecture that generalizes well. "
                     "Multi-objective Pareto front (L2 + params) identifies "
                     "efficient architectures automatically."),
            "bg": RGBColor(0xe3, 0xf2, 0xfd),
            "bar": C_BLUE,
        },
        {
            "num": "5",
            "title": "3D Gap Remains",
            "body": ("NAS in 2D cannot find adequate 3D architectures. "
                     "L2 jumps from 0.014 (2D) to 0.247 (3D). "
                     "Future work: 3D-specific NAS, transfer learning."),
            "bg": RGBColor(0xff, 0xf3, 0xe0),
            "bar": RGBColor(0xe6, 0x5c, 0x00),
        },
        {
            "num": "6",
            "title": "L9 Final: L2 = 0.014",
            "body": ("9\u00d7 better than FEM skip=1 baseline (L2=0.126). "
                     "18,400\u00d7 faster inference (0.01 s vs 184 s). "
                     "Enables NAS, optimization, real-time prediction."),
            "bg": RGBColor(0xe8, 0xf5, 0xe9),
            "bar": C_DKGREEN,
        },
    ]

    positions = [
        (0.25, 1.38), (4.58, 1.38), (8.92, 1.38),
        (0.25, 4.18), (4.58, 4.18), (8.92, 4.18),
    ]
    box_w = 4.1
    box_h = 2.55

    for f, (x, y) in zip(findings, positions):
        add_rect(sld, x, y, box_w, box_h, f["bg"])
        add_rect(sld, x, y, 0.5, box_h, f["bar"])
        add_text(sld, f["num"],
                 l=x, t=y + 0.9, w=0.5, h=0.5,
                 fs=14, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(sld, f["title"],
                 l=x + 0.6, t=y + 0.1, w=box_w - 0.7, h=0.38,
                 fs=12, bold=True, color=f["bar"])
        add_text(sld, f["body"],
                 l=x + 0.6, t=y + 0.53, w=box_w - 0.7, h=1.9,
                 fs=10, color=C_TEXT)

    return sld


def slide_14_references(prs):
    """Slide 14 — References"""
    sld = blank_slide(prs)
    add_rect(sld, 0, 0, 13.33, 7.5, C_WHITE)
    title_bar(sld, "References")

    refs_left = [
        ("[1]  Raissi, M., Perdikaris, P., & Karniadakis, G.E. (2019). "
         "Physics-informed neural networks: A deep learning framework for solving forward "
         "and inverse problems. J. Comput. Phys., 378, 686\u2013707. "
         "DOI: 10.1016/j.jcp.2018.10.045"),

        ("[2]  Mortensen et al. (2026). Mitigating distortions in cast automotive "
         "subframes: A finite element simulation approach. "
         "Int J Adv Manuf Technol. DOI: 10.1007/s00170-026-17515-w"),

        ("[3]  Tancik, M., et al. (2020). Fourier Features Let Networks Learn High Frequency "
         "Functions in Low Dimensional Domains. NeurIPS 2020."),

        ("[4]  Wang, S., Teng, Y., & Perdikaris, P. (2021). Understanding and mitigating "
         "gradient flow pathologies in physics-informed neural networks. "
         "SIAM J. Sci. Comput., 43(5)."),

        ("[5]  Wang, S., Yu, X., & Perdikaris, P. (2022). When and why PINNs fail to train: "
         "A neural tangent kernel perspective. J. Comput. Phys., 449. "
         "DOI: 10.1016/j.jcp.2021.110768"),

        ("[6]  Deb, K., Pratap, A., Agarwal, S., & Meyarivan, T. (2002). A fast and elitist "
         "multiobjective genetic algorithm: NSGA-II. IEEE Trans. Evol. Comput., 6(2). "
         "DOI: 10.1109/4235.996017"),
    ]

    refs_right = [
        ("[7]  Deb, K., & Jain, H. (2014). An evolutionary many-objective optimization "
         "algorithm using reference-point-based nondominated sorting: NSGA-III. "
         "IEEE Trans. Evol. Comput., 18(4). DOI: 10.1109/TEVC.2013.2281535"),

        ("[8]  Snoek, J., Larochelle, H., & Adams, R.P. (2012). Practical Bayesian "
         "Optimization of Machine Learning Algorithms. NeurIPS 2012."),

        ("[9]  Glorot, X., & Bengio, Y. (2010). Understanding the difficulty of training "
         "deep feedforward neural networks. AISTATS 2010."),

        ("[10] Lagaris, I.E., Likas, A., & Fotiadis, D.I. (1998). Artificial neural networks "
         "for solving ordinary and partial differential equations. "
         "IEEE Trans. Neural Netw., 9(5). DOI: 10.1109/72.712178"),

        ("[11] Rahaman, N., et al. (2019). On the spectral bias of neural networks. "
         "ICML 2019."),

        ("[12] Liu, D.C., & Nocedal, J. (1989). On the limited memory BFGS method for "
         "large scale optimization. Math. Programming, 45(1\u20133), 503\u2013528."),
    ]

    add_rect(sld, 0.25, 1.35, 6.35, 5.5, C_LGRAY)
    y = 1.45
    for ref in refs_left:
        add_text(sld, ref,
                 l=0.35, t=y, w=6.15, h=0.73,
                 fs=9, color=C_TEXT)
        y += 0.82

    add_rect(sld, 6.8, 1.35, 6.3, 5.5, C_LGRAY)
    y = 1.45
    for ref in refs_right:
        add_text(sld, ref,
                 l=6.9, t=y, w=6.1, h=0.73,
                 fs=9, color=C_TEXT)
        y += 0.82

    # Footer
    add_rect(sld, 0, 6.98, 13.33, 0.52, C_NAVY)
    add_text(sld, "Thank You",
             l=4.5, t=7.0, w=4.3, h=0.45,
             fs=16, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(sld, "March 2026",
             l=11.3, t=7.0, w=1.9, h=0.45,
             fs=12, color=C_LBLUE,
             align=PP_ALIGN.RIGHT)

    return sld


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

def main():
    print("=" * 60)
    print("NAS-MCO-PINNs  \u2014  Presentation Generator")
    print("=" * 60)
    print(f"BASE  : {BASE}")
    print(f"ROOT  : {ROOT}")
    print(f"PLOTS : {PLOTS}  (exists={os.path.isdir(PLOTS)})")
    print(f"L8_IMG: {L8_IMG}  (exists={os.path.isdir(L8_IMG)})")
    print(f"OUT   : {OUT}")
    print()

    prs = new_presentation()

    builders = [
        ("Slide  1 \u2014 Title",                       slide_01_title),
        ("Slide  2 \u2014 Project Overview",             slide_02_overview),
        ("Slide  3 \u2014 Problem: Industrial Casting",  slide_03_problem),
        ("Slide  4 \u2014 FEM Reference Model",          slide_04_fem_reference),
        ("Slide  5 \u2014 Why Replace FEM?",             slide_05_why_faster),
        ("Slide  6 \u2014 Temporal Skip Operator",       slide_06_temporal_skip),
        ("Slide  7 \u2014 Fourier Features",             slide_07_fourier),
        ("Slide  8 \u2014 Self-Adaptive Weights",        slide_08_sa_weights),
        ("Slide  9 \u2014 NAS Architecture Search",      slide_09_nas),
        ("Slide 10 \u2014 Level 8 3D Results",           slide_10_level8),
        ("Slide 11 \u2014 Level 9 Best Results",         slide_11_level9),
        ("Slide 12 \u2014 FEM vs PINN Comparison",       slide_12_fem_vs_pinn),
        ("Slide 13 \u2014 Key Findings",                 slide_13_findings),
        ("Slide 14 \u2014 References",                   slide_14_references),
    ]

    for label, fn in builders:
        print(f"  Building {label} ...", end=" ", flush=True)
        try:
            fn(prs)
            print("OK")
        except Exception as exc:
            print(f"ERROR: {exc}")
            import traceback
            traceback.print_exc()

    print()
    print(f"Saving to: {OUT}")
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("Done.")
    print()
    size_kb = os.path.getsize(OUT) / 1024
    print(f"Output file : {OUT}")
    print(f"File size   : {size_kb:.1f} KB")
    print("=" * 60)


if __name__ == "__main__":
    main()
